import os, sys, argparse, math
import numpy as np
from operator import itemgetter
from math import ceil
from datetime import date
from matplotlib.ticker import AutoMinorLocator
import matplotlib.pyplot as plt
import matplotlib
import matplotlib.ticker as ticker
from matplotlib.ticker import FormatStrFormatter
from matplotlib.ticker import ScalarFormatter
import matplotlib.gridspec as gridspec
import matplotlib.patches as patches
import periodictable
from colour import Color
import seaborn as sns
from scipy.stats import gaussian_kde as gkde
from scipy.stats import linregress
from scipy.interpolate import make_interp_spline, BSpline
import re

# For making PowerPoint Slides
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import ColorFormat, RGBColor

from astromol.molecules import *
from astromol.sources import *
from astromol.telescopes import *
import astromol

def version():
    return astromol.__version__

def updated():
    return f"{astromol.__updated__.day} {astromol.__updated__.strftime('%b')} {astromol.__updated__.year}"

matplotlib.rc("text", usetex=True)
matplotlib.rc("text.latex", preamble=r"\usepackage{cmbright}\usepackage[version=4]{mhchem}")

#############################################################
# 						Functions	 						#
#############################################################


def make_all_plots():
    """
    A meta function that, when run, will call every plot command and generate new plots based on
    the input list of Molecule objects using default parameters.  Useful for rapidly re-generating all figures.

    Will always use the full database and default filenames.
    """

    cumu_det_plot()
    cumu_det_natoms_plot()
    det_per_year_per_atom()
    facility_shares()
    scopes_by_year()
    periodic_heatmap()
    mass_by_wavelengths()
    mols_waves_by_atoms()
    du_histogram()
    type_pie_chart()
    source_pie_chart()
    indiv_source_pie_chart()
    mol_type_by_source_type()
    du_by_source_type()
    rel_du_by_source_type()
    mass_by_source_type()
    waves_by_source_type()

    return


def make_all_latex():
    """
    A meta function that, when run, will call every latex-generating command based on
    the input list of Molecule objects using default parameters.  Useful for rapidly re-generating all latex inputs.

    Will always use the full database and default filenames.
    """

    make_det_count()
    make_elem_count()
    make_ppd_count()
    make_ppd_isos_count()
    make_ism_tables()
    make_exgal_table()
    make_ppd_table()
    make_exo_table()
    make_ice_table()
    make_facility_table()
    make_source_table()
    make_rate_counts()
    make_exgal_count()
    make_exgal_percent()
    _make_exgal_sentence()
    make_exo_count()
    make_ices_count()
    make_percent_radio()
    make_scopes_count()
    make_percent_unsat()
    make_sat_list()
    make_sat_count()
    make_sat_percent()
    make_sfr_rad_percent()
    make_dark_rad_percent()
    make_det_per_year_by_atoms_table()

    return


def change_color(color, amount=1.0):
    """
    Lightens the given color by multiplying (1-luminosity) by the given amount.

    Parameters
    ----------
    color : str, tuple
        A matplotlib color string, hex string, or RGB tuple.
    amount : float
        The amount to lighten a color (default is 1.0)

    Returns
    -------
    color : tuple
        The altered color as an RGB tuple

    Examples
    --------
    >> lighten_color('g', 0.3)
    >> lighten_color('#F034A3', 0.6)
    >> lighten_color((.3,.55,.1), 0.5)
    """
    import matplotlib.colors as mc
    import colorsys

    try:
        c = mc.cnames[color]
    except:
        c = color
    c = colorsys.rgb_to_hls(*mc.to_rgb(c))
    return colorsys.hls_to_rgb(c[0], 1 - amount * (1 - c[1]), c[2])

def print_variables(type=None,natoms=None):
    """
    Prints out a list to the terminal of the variable names for molecules, telescopes, 
    and/or sources in the astromol database and the corresponding molecular formulas, names, 
    and so forth.

    The list is sorted by the number of atoms, and can be narrowed to choose only a certain
    number of atoms to print.

    Parameters
    ----------
    type: str
        Which type of item to print.  Choose from:
            "molecules"
            "telescopes"
            "sources"
        The default, None, will print all three
    natoms : int
        If specified, only molecules with the given number of atoms will be displayed 
        (default is None, which prints all molecules).
    """    

    if natoms is None:
        mol_list = all_molecules
    else:
        mol_list = [x for x in all_molecules if x.natoms == natoms]

    if natoms is not None and (type is None or type.lower() == 'molecules'):
        print("\n=======================================")
        print(f"            {natoms} Atom Molecules       ")
        print("=======================================\n")

        print(f"{'Variable Name':13} \t {'Formula':12} \t {'Name':35}")
        print("------------- \t ------------ \t -----------------------------------")

        for mol in mol_list:
            print(f'{mol.astromol_name:13} \t {mol.table_formula if mol.table_formula is not None else mol.formula:12} \t {mol.name:35}')

    elif type is None or type.lower() == 'molecules':
        for i in range(2,13):
            print("\n=======================================")
            print(f"            {i} Atom Molecules       ")
            print("=======================================\n")

            print(f"{'Variable Name':13} \t {'Formula':12} \t {'Name':35}")
            print("------------- \t ------------ \t -----------------------------------")

            for mol in [x for x in mol_list if x.natoms == i]:
                print(f'{mol.astromol_name:13} \t {mol.table_formula if mol.table_formula is not None else mol.formula:12} \t {mol.name:35}')

        print("\n=======================================")
        print(f"            PAH Molecules       ")
        print("=======================================\n")

        print(f"{'Variable Name':13} \t {'Formula':12} \t {'Name':35}")
        print("------------- \t ------------ \t -----------------------------------")

        for mol in [x for x in mol_list if x.pah is True]:
            print(f'{mol.astromol_name:13} \t {mol.table_formula if mol.table_formula is not None else mol.formula:12} \t {mol.name:35}')

        print("\n=======================================")
        print(f"           Fullerene Molecules       ")
        print("=======================================\n")

        print(f"{'Variable Name':13} \t {'Formula':12} \t {'Name':35}")
        print("------------- \t ------------ \t -----------------------------------")

        for mol in [x for x in mol_list if x.fullerene is True]:
            print(f'{mol.astromol_name:13} \t {mol.table_formula if mol.table_formula is not None else mol.formula:12} \t {mol.name:35}')

    if (type is None or type.lower() == 'telescopes') and natoms is None:
        print("\n=======================================")
        print(f"              Telescopes       ")
        print("=======================================\n")

        print(f"{'Variable Name':13} \t {'Telescope Name':35}")
        print("------------- \t -----------------------------------")  

        for scope in all_telescopes:
            print(f'{scope.astromol_name:13} \t {scope.name:35}')     

    if (type is None or type.lower() == 'sources') and natoms is None:
        print("\n=======================================")
        print(f"              Sources       ")
        print("=======================================\n")

        print(f"{'Variable Name':13} \t {'Source Name':35}")
        print("------------- \t -----------------------------------")  

        for source in all_sources:
            print(f'{source.astromol_name:13} \t {source.name:35}')


def inspect(x):
    """
    Prints out a list to the terminal all of the attributes for the input item and their values.

    Note that this will be a largely unformatted output, and is really only intended to be used
    when a detailed look at what information is in the database for a given species is desired.

    Further note that this is just calling the built-in class method 'inspect'.

    Parameters
    ----------
    x : Molecule, Telescope, Source
        The molecule, telescope, or source object to inspect and print out
    """    

    x.inspect()

def summary(molecule):
    """
    Prints a nicely formatted list of references and notes to the terminal.

    Note that this is just calling the built-in Molecule class method 'summary'.

    Parameters
    ----------
    molecule : Molecule
        The molecule object to inspect and print out
    """    

    molecule.summary()

#############################################################
# 						    Plots	 						#
#############################################################


def cumu_det_plot(mol_list=None, syear=None, eyear=None, filename=None):

    """
    Makes a plot of the cumulative detections by year.

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)
    syear : int
        The starting year to plot (default is the earliest year of a molecule 
        detection from molecules in mol_list)
    eyear: int
        The ending year to plot (default is the current year).  Actual plot maximum
        will be one higher than this number for readability.
    filename : str
        The filename for the output images (default is 'cumulative_detections.pdf')
    """

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    # Close an old figure if it exists and initialize a new figure
    plt.close("Cumulative Detections")
    fig = plt.figure(num="Cumulative Detections", figsize=(10, 8))
    plt.ion()

    # set some font defaults
    fontparams = {"size": 24, "family": "sans-serif", "sans-serif": ["Helvetica"]}
    plt.rc("font", **fontparams)
    plt.rc("mathtext", fontset="stixsans")

    # get the starting and ending years, if they aren't set by the user
    if syear is None:
        # grab the earliest year in the list of molecules
        syear = min([x.year for x in mol_list])

    if eyear is None:
        # grab the current year
        eyear = date.today().year

    # make the x-axis array of years
    years = np.arange(syear, eyear + 1)

    # make an array to hold the detections
    dets = np.copy(years) * 0

    # loop through the years and the list and add everything up.  There's gotta be a better way to do this, but who cares, its fast enough.
    for x in range(len(dets)):
        i = 0
        for mol in mol_list:
            if mol.year < years[x] + 1:
                i += 1
        dets[x] = i

    # get some year indicies for years we care about
    def iyear(x):
        return np.argwhere(years == x)[0][0]

    # do linear fits to the data for the two ranges we care about (1968-Present, 2005-Present)
    # get the slope of the fit to detections since 1968 and since 2005

    trend1968 = np.polynomial.polynomial.Polynomial.fit(years[iyear(1968) :], dets[iyear(1968) :], 1).convert().coef[1]
    trend2005 = np.polynomial.polynomial.Polynomial.fit(years[iyear(2005) :], dets[iyear(2005) :], 1).convert().coef[1]

    # load up an axis
    ax = fig.add_subplot(111)

    # label the axes
    plt.xlabel("Year")
    plt.ylabel("Cumulative Number of Detected Molecules")

    # customize tick marks
    ax.tick_params(axis="x", which="both", direction="in", length=15, width=1)
    ax.tick_params(axis="y", which="both", direction="in", length=15, width=1)

    ax.yaxis.set_ticks_position("both")
    ax.xaxis.set_ticks_position("both")

    # plot it
    ax.plot(years, dets, color="dodgerblue", linewidth=4)

    # add annotations; #first, the detections per year
    args = {"ha": "left", "size": "24"}
    det_str = r"\noindent Since 1968: {:.1f} detections/year \\ Since 2005: {:.1f} detections/year".format(
        trend1968, trend2005
    )
    ax.annotate(det_str, xy=(0.05, 0.85), xycoords="axes fraction", **args)

    # Now the total number of detections
    ax.annotate(
        "Total: {}".format(dets[iyear(eyear)]),
        xy = (0.93,0.93),
        xycoords = 'axes fraction',
        va = 'top',
        ha = 'right',
        size=24,
    )

    # Now the facilities
    args = {"size": "16"}
    arrowprops = {"arrowstyle": "-|>", "connectionstyle": "arc3", "facecolor": "black"}
    ax.annotate(
        "NRAO 36-foot (1968)",
        xy=(1968, dets[iyear(1968)] + 4),
        xycoords="data",
        xytext=(0, 35),
        textcoords="offset points",
        rotation=90,
        arrowprops=arrowprops,
        va="bottom",
        ha="center",
        **args,
    )
    ax.annotate(
        "",
        xy=(1982, dets[iyear(1982)] - 4),
        xycoords="data",
        xytext=(0, -35),
        textcoords="offset points",
        rotation=0,
        arrowprops=arrowprops,
        va="bottom",
        ha="left",
        **args,
    )
    ax.annotate("Nobeyama (1982)", xy=(1982, dets[iyear(1982)] - 31), xycoords="data", **args)
    ax.annotate(
        "IRAM (1984)",
        xy=(1984, dets[iyear(1984)] + 4),
        xycoords="data",
        xytext=(0, 35),
        textcoords="offset points",
        rotation=90,
        arrowprops=arrowprops,
        va="bottom",
        ha="center",
        **args,
    )
    ax.annotate(
        "GBT (2001)",
        xy=(2001, dets[iyear(2001)] + 4),
        xycoords="data",
        xytext=(0, 35),
        textcoords="offset points",
        rotation=90,
        arrowprops=arrowprops,
        va="bottom",
        ha="center",
        **args,
    )
    ax.annotate(
        "ALMA (2011)",
        xy=(2011, dets[iyear(2011)] - 4),
        xycoords="data",
        xytext=(0, -35),
        textcoords="offset points",
        rotation=90,
        arrowprops=arrowprops,
        va="top",
        ha="center",
        **args,
    )
    ax.annotate(
        "Yebes (2007)",
        xy=(2007, dets[iyear(2007)] - 4),
        xycoords="data",
        xytext=(0, -35),
        textcoords="offset points",
        rotation=90,
        arrowprops=arrowprops,
        va="top",
        ha="center",
        **args,
    )    

    # write out the figure
    plt.savefig(
        filename if filename is not None else "cumulative_detections.pdf",
        format="pdf",
        transparent=True,
        bbox_inches="tight",
    )
    
    # show the plot
    plt.show()


def cumu_det_natoms_plot(mol_list=None, syear=None, eyear=None, filename=None):
    """
    Makes a plot of the cumulative detections (sorted by atoms) by year

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)
    syear : int
        The starting year to plot (default is the earliest year of a molecule 
        detection from molecules in mol_list)
    eyear: int
        The ending year to plot (default is the current year).  Actual plot maximum
        will be one year higher than this number for readability.
    filename : str
        The filename for the output images (default is 'cumulative_by_atoms.pdf')
    """

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    # Close an old figure if it exists and initialize a new figure
    plt.close("Cumulative Detections By Atoms")
    fig = plt.figure(num="Cumulative Detections By Atoms", figsize=(10, 8))
    plt.ion()

    # set some font defaults
    fontparams = {"size": 24, "family": "sans-serif", "sans-serif": ["Helvetica"]}
    plt.rc("font", **fontparams)
    plt.rc("mathtext", fontset="stixsans")

    # get the starting and ending years, if they aren't set by the user

    if syear is None:
        # grab the earliest year in the list of molecules
        syear = min([x.year for x in mol_list])

    if eyear is None:
        # grab the current year
        eyear = date.today().year

    # make the x-axis array of years
    years = np.arange(syear, eyear + 1)

    # now we make a dictionary of the different traces we're gonna want, loop through the list, and populate an array of detections for that number of atoms or special case
    dets_dict = {}
    for natoms in range(2, 14):
        # fill it with an empte. array
        dets_dict[natoms] = np.copy(years) * 0
        # loop through the years and the list and add everything up.  There's gotta be a better way to do this, but who cares, its fast enough.
        for x in range(len(dets_dict[natoms])):
            i = 0
            for mol in mol_list:
                if mol.year < years[x] + 1 and mol.natoms == natoms:
                    i += 1
            dets_dict[natoms][x] = i

    # do the fullerenes and pahs
    dets_dict["fullerenes"] = np.copy(years) * 0
    for x in range(len(dets_dict["fullerenes"])):
        i = 0
        for mol in mol_list:
            if mol.year < years[x] + 1 and mol.fullerene is True:
                i += 1
        dets_dict["fullerenes"][x] = i

    dets_dict["pahs"] = np.copy(years) * 0
    for x in range(len(dets_dict["pahs"])):
        i = 0
        for mol in mol_list:
            if mol.year < years[x] + 1 and mol.pah is True:
                i += 1
        dets_dict["pahs"][x] = i

    # load up an axis
    ax = fig.add_subplot(111)

    # label the axes
    plt.xlabel("Year")
    plt.ylabel("Cumulative Number of Detected Molecules")

    # customize tick marks
    ax.tick_params(axis="x", which="both", direction="in", length=15, width=1)
    ax.tick_params(axis="y", which="both", direction="in", length=15, width=1)
    ax.yaxis.set_ticks_position("both")
    ax.xaxis.set_ticks_position("both")
    ax.set_xticks([1940, 1960, 1980, 2000, 2020])
    ax.set_xlim([syear, eyear])

    # plot the traces
    ax.plot(years, dets_dict[2], color="#000000")
    ax.plot(years, dets_dict[3], color="#800000")
    ax.plot(years, dets_dict[4], color="#f032e6")
    ax.plot(years, dets_dict[5], color="#9A6324")
    ax.plot(years, dets_dict[6], color="dodgerblue")
    ax.plot(years, dets_dict[7], color="#e6194B")
    ax.plot(years, dets_dict[8], color="#469990")
    ax.plot(years, dets_dict[9], color="#f58231")
    ax.plot(years, dets_dict[10], color="#42d4f4")
    ax.plot(years, dets_dict[11], color="#ffe119")
    ax.plot(years, dets_dict[12], color="#3cb44b")
    ax.plot(years, dets_dict[13], color="#e6beff")
    ax.plot(years, dets_dict["fullerenes"], color="#000075")
    ax.plot(years, dets_dict["pahs"], color="#aaffc3")

    # add the annotations

    start_y = 0.92
    step_y = 0.04

    start_x = 0.10

    ax.annotate(
        r"\textbf{2 atoms}",
        xy=(start_x, start_y - (0 * step_y)),
        xycoords="axes fraction",
        size=16,
        color="#000000",
        va="top",
        ha="left",
    )
    ax.annotate(
        r"\textbf{3 atoms}",
        xy=(start_x, start_y - (1 * step_y)),
        xycoords="axes fraction",
        size=16,
        color="#800000",
        va="top",
        ha="left",
    )
    ax.annotate(
        r"\textbf{4 atoms}",
        xy=(start_x, start_y - (2 * step_y)),
        xycoords="axes fraction",
        size=16,
        color="#f032e6",
        va="top",
        ha="left",
    )
    ax.annotate(
        r"\textbf{5 atoms}",
        xy=(start_x, start_y - (3 * step_y)),
        xycoords="axes fraction",
        size=16,
        color="#9A6324",
        va="top",
        ha="left",
    )
    ax.annotate(
        r"\textbf{6 atoms}",
        xy=(start_x, start_y - (4 * step_y)),
        xycoords="axes fraction",
        size=16,
        color="dodgerblue",
        va="top",
        ha="left",
    )
    ax.annotate(
        r"\textbf{7 atoms}",
        xy=(start_x, start_y - (5 * step_y)),
        xycoords="axes fraction",
        size=16,
        color="#e6194B",
        va="top",
        ha="left",
    )
    ax.annotate(
        r"\textbf{8 atoms}",
        xy=(start_x, start_y - (6 * step_y)),
        xycoords="axes fraction",
        size=16,
        color="#469990",
        va="top",
        ha="left",
    )
    ax.annotate(
        r"\textbf{9 atoms}",
        xy=(start_x, start_y - (7 * step_y)),
        xycoords="axes fraction",
        size=16,
        color="#f58231",
        va="top",
        ha="left",
    )
    ax.annotate(
        r"\textbf{10 atoms}",
        xy=(start_x, start_y - (8 * step_y)),
        xycoords="axes fraction",
        size=16,
        color="#42d4f4",
        va="top",
        ha="left",
    )
    ax.annotate(
        r"\textbf{11 atoms}",
        xy=(start_x, start_y - (9 * step_y)),
        xycoords="axes fraction",
        size=16,
        color="#ffe119",
        va="top",
        ha="left",
    )
    ax.annotate(
        r"\textbf{12 atoms}",
        xy=(start_x, start_y - (10 * step_y)),
        xycoords="axes fraction",
        size=16,
        color="#3cb44b",
        va="top",
        ha="left",
    )
    ax.annotate(
        r"\textbf{13+ atoms}",
        xy=(start_x, start_y - (11 * step_y)),
        xycoords="axes fraction",
        size=16,
        color="#e6beff",
        va="top",
        ha="left",
    )

    ax.annotate(
        r"\textbf{Fullerenes}",
        xy=(start_x, start_y - (12 * step_y)),
        xycoords="axes fraction",
        size=16,
        color="#000075",
        va="top",
        ha="left",
    )
    ax.annotate(
        r"\textbf{PAHs}",
        xy=(start_x, start_y - (13 * step_y)),
        xycoords="axes fraction",
        size=16,
        color="#aaffc3",
        va="top",
        ha="left",
    )

    # show the plot
    plt.show()

    # write out the figure
    plt.savefig(
        filename if filename is not None else "cumulative_by_atoms.pdf",
        format="pdf",
        transparent=True,
        bbox_inches="tight",
    )


def det_per_year_per_atom(mol_list=None, filename=None):
    """
    Makes a plot of the average number of detections per year (y) for a molecule with (x) atoms, 
    starting in the year they were first detected.

    Has the ability to plot PAHs and fullerenes, but doesn't.

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)
    filename : str
        The filename for the output images (default is 'rate_by_atoms.pdf')
    """

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    # Close an old figure if it exists and initialize a new figure
    plt.close("Detects Per Year Per Atom")
    fig = plt.figure(num="Detects Per Year Per Atom", figsize=(10, 8))
    plt.ion()

    # set some font defaults
    fontparams = {"size": 24, "family": "sans-serif", "sans-serif": ["Helvetica"]}
    plt.rc("font", **fontparams)
    plt.rc("mathtext", fontset="stixsans")

    # We're going to cheat later on and lump PAHs and Fullerenes together, but these need fake natoms.
    # They'll always be +2 and +4, respectively, beyond the maximum of 'normal' molecules, and we'll
    # leave +1 blank for visual separation.  For now, though, time to make a list and loop through our molecules.
    natoms = np.arange(2, 18)

    # and an array of the average number of detections
    avg_dets = np.copy(natoms) * 0.0
    n_dets = np.copy(natoms) * 0.0

    # get the years that are being spanned
    eyear = max([x.year for x in mol_list])
    for x in range(len(natoms)):
        i = 0
        years = []
        if natoms[x] < 14:
            for mol in mol_list:
                if mol.natoms == natoms[x]:
                    i += 1
                    years.append(mol.year)
        elif natoms[x] == 15:
            for mol in mol_list:
                if mol.pah is True:
                    i += 1
                    years.append(mol.year)
        elif natoms[x] == 17:
            for mol in mol_list:
                if mol.fullerene is True:
                    i += 1
                    years.append(mol.year)
        if i == 0:
            avg_dets[x] = np.NaN
        else:
            avg_dets[x] = i / (eyear - min(years) + 1)
            n_dets[x] = i
    n_dets[n_dets == 0] = np.nan

    # load up an axis
    ax = fig.add_subplot(111)

    # label the axes
    plt.xlabel("Number of Atoms")
    plt.ylabel("Detections/Year*")

    # customize tick marks
    ax.tick_params(axis="x", which="both", direction="in", length=15, width=1)
    ax.tick_params(axis="y", which="both", direction="in", length=15, width=1)
    ax.yaxis.set_ticks_position("both")
    ax.xaxis.set_ticks_position("both")
    ax.set_xticks([2, 4, 6, 8, 10, 12, 15, 17])
    ax.set_xticklabels(["2", "4", "6", "8", "10", "12", "PAHs", "Fullerenes"])
    ax.set_xlim([1, 12.8])
    ax.set_ylim([0, 1.0])

    # plot the data
    mfc = "dodgerblue"
    mec = change_color(mfc, 1.5)
    sizes = np.copy(n_dets) * 200 / np.nanmin(n_dets)
    ax.scatter(natoms, avg_dets, marker="o", c=mfc, edgecolors=mec, s=sizes, alpha=0.9)

    # add some annotations
    ax.annotate(
        "*Since year of first detection",
        xy=(0.35, 0.9),
        xycoords="axes fraction",
        ha="left",
    )
    ax.annotate(
        r"\noindent Marker size proportional to\\ total \# of detections",
        xy=(0.05, 0.1),
        xycoords="axes fraction",
        ha="left",
    )
    for label in ax.get_xmajorticklabels():
        if label._text == "Fullerenes" or label._text == "PAHs":
            label.set_rotation(-45)
    fig.tight_layout()

    # show the plot
    plt.show()

    # write out the figure
    plt.savefig(
        filename if filename is not None else "rate_by_atoms.pdf",
        format="pdf",
        transparent=True,
        bbox_inches="tight",
    )

    return


def facility_shares(mol_list=None, telescopes_list=None, filename=None):
    """
    Makes a plot of the percentage share of yearly detections that a facility 
    contributed over its operational lifetime for the top 9 facilities.

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)
    telescopes_list : list
        A list of telescope objects to use (default is all_telescopes)
    filename : str
        The filename for the output images (default is 'facility_shares.pdf')
    """

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules
    # If a list wasn't specified, default to all telescopes
    if telescopes_list is None:
        telescopes_list = all_telescopes

    # Close an old figure if it exists and initialize a new figure
    plt.close("Facility Shares")
    fig, axs = plt.subplots(3, 3, num="Facility Shares")
    plt.ion()

    # set some font defaults
    fontparams = {"size": 14, "family": "sans-serif", "sans-serif": ["Helvetica"]}
    plt.rc("font", **fontparams)
    plt.rc("mathtext", fontset="stixsans")

    # we need to generate the data now, which we'll store in a dictionary for each telescope.
    # Each entry will be [syear,eyear,ndetects,ntotal,shortname] for the start year, end year,
    # number of detections, and total number of detections over those years
    detects = [x.ndetects for x in telescopes_list]
    detects.sort(reverse=True)
    # determine the 9 most productive facilities on an absolute scale
    detects = detects[:9]
    # get the minimum number of detections of these
    min_allowed = min(detects)

    my_dict = {}
    for scope in telescopes_list:
        # number of detections
        ndetects = scope.ndetects
        # determine if this telescope has enough detections to qualify
        # if it does, just pass and do the rest
        if ndetects >= min_allowed:
            pass
        # otherwise, move on to the next
        else:
            continue
        # year it was built
        syear = scope.built
        # year it was decommissioned or this year if it's still in operation
        eyear = scope.decommissioned if scope.decommissioned is not None else date.today().year

        # now we go get the total number of detections in that time
        ntotal = 0
        for mol in mol_list:
            if mol.year <= eyear and mol.year >= syear:
                ntotal += 1
        my_dict[scope.shortname] = [syear, eyear, ndetects, ntotal, scope.shortname]

    my_list = [my_dict[x] for x in my_dict]
    my_list.sort(key=lambda x: x[2] / x[3], reverse=True)

    idx = 0
    for i in range(3):
        for j in range(3)[::]:
            fracs = [
                my_list[idx][2] / my_list[idx][3],
                1.0 - my_list[idx][2] / my_list[idx][3],
            ]
            label = r"\textbf" + "{" + "{}".format(int(my_list[idx][2] / my_list[idx][3] * 100)) + r"\%"

            if my_list[idx][1] == date.today().year:
                color = "dodgerblue"
            else:
                color = "#F87070"

            if idx > 5:
                slices, labels = axs[i, j].pie(
                    fracs,
                    colors=[color, "#F5F6FF"],
                    wedgeprops={"linewidth": 1.0, "edgecolor": "black"},
                )

                axs[i, j].annotate(
                    my_list[idx][-1],
                    xy=(0.5, 1.08),
                    xycoords="axes fraction",
                    ha="center",
                    size=12,
                )
                axs[i, j].annotate(
                    "{} - {}".format(my_list[idx][0], my_list[idx][1]),
                    xy=(0.5, 0.95),
                    xycoords="axes fraction",
                    ha="center",
                    size=10,
                )

                kw = dict(arrowprops=dict(arrowstyle="-"), zorder=0, va="center")
                ang = (slices[0].theta2 - slices[0].theta1) / 2.0 + slices[0].theta1
                y = np.sin(np.deg2rad(ang))
                x = np.cos(np.deg2rad(ang))
                horizontalalignment = {-1: "right", 1: "left"}[int(np.sign(x))]
                connectionstyle = "angle,angleA=0,angleB={}".format(ang)
                kw["arrowprops"].update({"connectionstyle": connectionstyle})
                axs[i, j].annotate(
                    label,
                    xy=(x, y),
                    xytext=(1.1 * np.sign(x), y),
                    horizontalalignment=horizontalalignment,
                    size=10,
                    **kw,
                )

            else:
                my_labels = [
                    r"\textbf" + "{" + "{}".format(int(my_list[idx][2] / my_list[idx][3] * 100)) + r"\%",
                    "",
                ]

                slices, labels = axs[i, j].pie(
                    fracs,
                    labels=my_labels,
                    colors=[color, "#F5F6FF"],
                    wedgeprops={"linewidth": 1.0, "edgecolor": "black"},
                    labeldistance=0.7,
                )

                axs[i, j].annotate(
                    my_list[idx][-1],
                    xy=(0.5, 1.08),
                    xycoords="axes fraction",
                    ha="center",
                    size=12,
                )
                axs[i, j].annotate(
                    "{} - {}".format(my_list[idx][0], my_list[idx][1]),
                    xy=(0.5, 0.95),
                    xycoords="axes fraction",
                    ha="center",
                    size=10,
                )

                for label in labels:
                    label.set_horizontalalignment("center")
                    label.set_color("white")
                    label.set_fontsize(10)

            # draw a line around the outside of the pie
            # get the center and radius of the pie wedges
            center = slices[0].center
            r = slices[0].r
            # create a new circle with the desired properties
            circle = matplotlib.patches.Circle(center, r, fill=False, edgecolor="black", linewidth=1)
            # add the circle to the axes
            axs[i, j].add_patch(circle)

            idx += 1

    fig.tight_layout()
    fig.subplots_adjust(wspace=-0.5, hspace=0.15)
    plt.show()

    plt.savefig(
        filename if filename is not None else "facility_shares.pdf",
        format="pdf",
        transparent=True,
        bbox_inches="tight",
    )


def scopes_by_year(mol_list=None, telescopes_list=None, min_detects=10, filename=None):
    """
    Makes a plot of the cumulative number of detections of a facility with time.

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)
    telescopes_list : list
        A list of telescope objects to use (default is all_telescopes)
    min_detects : int
        The minimum number of detections required to include a facility in the plot
        (default is 10)
    filename : str
        The filename for the output images (default is 'scopes_by_year.pdf')
    """

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules
    # If a list wasn't specified, default to all telescopes
    if telescopes_list is None:
        telescopes_list = all_telescopes

    # Close an old figure if it exists and initialize a new figure
    plt.close("Detections Per Facility Over Time")
    fig = plt.figure(num="Detections Per Facility Over Time", figsize=(10, 8))
    plt.ion()

    # set some font defaults
    fontparams = {"size": 24, "family": "sans-serif", "sans-serif": ["Helvetica"]}
    plt.rc("font", **fontparams)
    plt.rc("mathtext", fontset="stixsans")

    # We're only going to do facilities with min_detects or more total detections.
    # set this thing up to kick off a few years before 1968 and run until today

    years = np.arange(1965, date.today().year + 1)
    scopes = [x for x in telescopes_list if x.ndetects >= min_detects]
    # need a color pallette to work with.  This one is supposedly color-blind friendly.  There are more here than currently needed.
    # if we happen to exceed this at some point, it will crash.
    colors = ["#377eb8", "#ff7f00", "#4daf4a", "#f781bf", "#a65628", "#984ea3", "#999999", "#e41a1c", "#dede00"][
        : len(scopes)
    ]

    my_dict = {}
    for scope, color in zip(scopes, colors):
        tmp_years = np.copy(years) * 0
        i = 0
        for x in range(len(years)):
            for mol in mol_list:
                if mol.year == years[x] and scope in mol.telescopes:
                    i += 1
            tmp_years[x] = i
        my_dict[scope.shortname] = [tmp_years, color, scope]

    # do linear fits to the data for the ranges we care about for each facility:
    # get some year indicies for years we care about
    def iyear(x):
        return np.argwhere(years == x)[0][0]

    # manually add cutoff dates for telescopes to be used in the fitting.  If a date isn't added, it defaults to the present
    cutoffs = {
        "NRAO 36-ft": 1985,
        "NRAO 140-ft": 1993,
        "Nobeyama 45-m": 1997,
    }

    # add trends to the dictionary
    for x in my_dict:
        end_year = iyear(cutoffs[my_dict[x][2].shortname]) if my_dict[x][2].shortname in cutoffs else None
        my_dict[x].append(
            np.polynomial.polynomial.Polynomial.fit(
                years[iyear(my_dict[x][2].built) : end_year], my_dict[x][0][iyear(my_dict[x][2].built) : end_year], 1
            )
            .convert()
            .coef[1]
        )

    # pull up an axis
    ax = fig.add_subplot(111)

    # label the axes
    plt.xlabel("Year")
    plt.ylabel("Cumulative Number of Detected Molecules")

    # customize tick marks
    ax.tick_params(axis="x", which="both", direction="in", length=15, width=1)
    ax.tick_params(axis="y", which="both", direction="in", length=15, width=1)
    ax.yaxis.set_ticks_position("both")
    ax.xaxis.set_ticks_position("both")
    ax.set_xticks([1970, 1980, 1990, 2000, 2010, 2020])

    # plot
    for x in my_dict:
        ax.plot(years, my_dict[x][0], color=my_dict[x][1])

    # we want to sort the labels with the most productive telescopes on top; that's hard to do with a dictionary
    # so, we get the rates into a list, and the keys into a list, then sort both of those by the rates, then iterate over the sorted keys
    rates = []
    keys = []
    for x in my_dict:
        keys.append(x)
        rates.append(my_dict[x][3])
    srt_idx = np.argsort(np.array(rates))[::-1]
    keys = np.array(keys)[srt_idx]

    row = 0
    for x in keys:
        ax.annotate(
            my_dict[x][2].shortname,
            xy=(0.1, 0.90 - 0.05 * row),
            xycoords="axes fraction",
            color=my_dict[x][1],
            fontweight="bold",  # for some reason, matplotlib is entirely ignoring this
            ha="left",
            va="top",
            size=18,
        )
        ax.annotate(
            f"{my_dict[x][3]:.1f}/yr",
            xy=(0.37, 0.90 - 0.05 * row),
            xycoords="axes fraction",
            color=my_dict[x][1],
            fontweight="bold",  # for some reason, matplotlib is entirely ignoring this
            ha="left",
            va="top",
            size=18,
        )
        date_str_a = f"({my_dict[x][2].built} -"
        date_str_b = (
            f"{cutoffs[my_dict[x][2].shortname]})" if my_dict[x][2].shortname in cutoffs else f"{date.today().year})"
        )
        ax.annotate(
            date_str_a + date_str_b,
            xy=(0.47, 0.90 - 0.05 * row),
            xycoords="axes fraction",
            color=my_dict[x][1],
            fontweight="bold",  # for some reason, matplotlib is entirely ignoring this
            ha="left",
            va="top",
            size=18,
        )
        row += 1

    ax.set_xlim([1965, date.today().year + 2])
    ax.set_ylim(0, max([my_dict[x][0][-1] for x in my_dict]) + 5)

    plt.show()

    plt.savefig(
        filename if filename is not None else "scopes_by_year.pdf",
        format="pdf",
        transparent=True,
        bbox_inches="tight",
    )


def periodic_heatmap(mol_list=None, filename=None, pdf_crop=True):
    """
    Makes a periodic table heat map of molecules detected containing each element.

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)
    filename : str
        The filename for the output images (default is 'periodic_heatmap.pdf')
    pdf_crop : bool
        If True, will use the system's TexLive pdfcrop utility to crop off extra
        whitespace from the final PDF.  Likely only works on Macs (default is True)
    """

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    # Close an old figure if it exists and initialize a new figure
    plt.close("Periodic Heatmap")
    plt.figure(num="Periodic Heatmap", figsize=(20, 9.5))
    plt.ion()

    # set some font defaults
    fontparams = {"size": 24, "family": "sans-serif", "sans-serif": ["Helvetica"]}
    plt.rc("font", **fontparams)
    plt.rc("mathtext", fontset="stixsans")
    # test
    # get a list of elements involved - we'll use the masses list from molecules.py, as that will be update to date
    els = list(masses.keys())
    # open a new dictionary to store detections from these
    census = {}

    # loop through, creating the dictionary entry if needed
    for mol in mol_list:
        for el in els:
            if el in mol.atoms:
                if mol.atoms[el] > 0:
                    if el in census:
                        census[el] += 1
                    else:
                        census[el] = 1

    maxdets = max([census[x] for x in census])
    map_colors = list(Color("#f1fb53").range_to(Color("#f00707"), maxdets))
    map_colors = [str(x) for x in map_colors]

    # Dictionary for the periodic table
    elements = {
        "H": [periodictable.H, 1, 6.9],
        "He": [periodictable.He, 18, 6.9],
        "Li": [periodictable.Li, 1, 5.75],
        "Be": [periodictable.Be, 2, 5.75],
        "B": [periodictable.B, 13, 5.75],
        "C": [periodictable.C, 14, 5.75],
        "N": [periodictable.N, 15, 5.75],
        "O": [periodictable.O, 16, 5.75],
        "F": [periodictable.F, 17, 5.75],
        "Ne": [periodictable.Ne, 18, 5.75],
        "Na": [periodictable.Na, 1, 4.6],
        "Mg": [periodictable.Mg, 2, 4.6],
        "Al": [periodictable.Al, 13, 4.6],
        "Si": [periodictable.Si, 14, 4.6],
        "P": [periodictable.P, 15, 4.6],
        "S": [periodictable.S, 16, 4.6],
        "Cl": [periodictable.Cl, 17, 4.6],
        "Ar": [periodictable.Ar, 18, 4.6],
        "K": [periodictable.K, 1, 3.45],
        "Ca": [periodictable.Ca, 2, 3.45],
        "Sc": [periodictable.Sc, 3, 3.45],
        "Ti": [periodictable.Ti, 4, 3.45],
        "V": [periodictable.V, 5, 3.45],
        "Cr": [periodictable.Cr, 6, 3.45],
        "Mn": [periodictable.Mn, 7, 3.45],
        "Fe": [periodictable.Fe, 8, 3.45],
        "Co": [periodictable.Co, 9, 3.45],
        "Ni": [periodictable.Ni, 10, 3.45],
        "Cu": [periodictable.Cu, 11, 3.45],
        "Zn": [periodictable.Zn, 12, 3.45],
        "Ga": [periodictable.Ga, 13, 3.45],
        "Ge": [periodictable.Ge, 14, 3.45],
        "As": [periodictable.As, 15, 3.45],
        "Se": [periodictable.Se, 16, 3.45],
        "Br": [periodictable.Br, 17, 3.45],
        "Kr": [periodictable.Kr, 18, 3.45],
        "Rb": [periodictable.Rb, 1, 2.3],
        "Sr": [periodictable.Sr, 2, 2.3],
        "Y": [periodictable.Y, 3, 2.3],
        "Zr": [periodictable.Zr, 4, 2.3],
        "Nb": [periodictable.Nb, 5, 2.3],
        "Mo": [periodictable.Mo, 6, 2.3],
        "Tc": [periodictable.Tc, 7, 2.3],
        "Ru": [periodictable.Ru, 8, 2.3],
        "Rh": [periodictable.Rh, 9, 2.3],
        "Pd": [periodictable.Pd, 10, 2.3],
        "Ag": [periodictable.Ag, 11, 2.3],
        "Cd": [periodictable.Cd, 12, 2.3],
        "In": [periodictable.In, 13, 2.3],
        "Sn": [periodictable.Sn, 14, 2.3],
        "Sb": [periodictable.Sb, 15, 2.3],
        "Te": [periodictable.Te, 16, 2.3],
        "I": [periodictable.I, 17, 2.3],
        "Xe": [periodictable.Xe, 18, 2.3],
        "Cs": [periodictable.Cs, 1, 1.15],
        "Ba": [periodictable.Ba, 2, 1.15],
        "Hf": [periodictable.Hf, 4, 1.15],
        "Ta": [periodictable.Ta, 5, 1.15],
        "W": [periodictable.W, 6, 1.15],
        "Re": [periodictable.Re, 7, 1.15],
        "Os": [periodictable.Os, 8, 1.15],
        "Ir": [periodictable.Ir, 9, 1.15],
        "Pt": [periodictable.Pt, 10, 1.15],
        "Au": [periodictable.Au, 11, 1.15],
        "Hg": [periodictable.Hg, 12, 1.15],
        "Tl": [periodictable.Tl, 13, 1.15],
        "Pb": [periodictable.Pb, 14, 1.15],
        "Bi": [periodictable.Bi, 15, 1.15],
        "Po": [periodictable.Po, 16, 1.15],
        "At": [periodictable.At, 17, 1.15],
        "Rn": [periodictable.Rn, 18, 1.15],
        "Fr": [periodictable.Fr, 1, 0.0],
        "Ra": [periodictable.Ra, 2, 0.0],
        "Rf": [periodictable.Rf, 4, 0.0],
        "Db": [periodictable.Db, 5, 0.0],
        "Sg": [periodictable.Sg, 6, 0.0],
        "Bh": [periodictable.Bh, 7, 0.0],
        "Hs": [periodictable.Hs, 8, 0.0],
        "Mt": [periodictable.Mt, 9, 0.0],
        "Ds": [periodictable.Ds, 10, 0.0],
        "Rg": [periodictable.Rg, 11, 0.0],
        "Cn": [periodictable.Cn, 12, 0.0],
        "Nh": [periodictable.Nh, 13, 0.0],
        "Fl": [periodictable.Fl, 14, 0.0],
        "Mc": [periodictable.Mc, 15, 0.0],
        "Lv": [periodictable.Lv, 16, 0.0],
        "Ts": [periodictable.Ts, 17, 0.0],
        "Og": [periodictable.Og, 18, 0.0],
    }

    # load up an axis
    ax = plt.axes([0, 0, 1, 1])
    ax.set_xlim([0, 18])
    ax.set_ylim([0, 8])

    for el in elements:
        x = elements[el][1] - 1
        y = elements[el][2]
        sym = r"\textbf{" + elements[el][0].symbol + "}"
        num = elements[el][0].number
        mass = elements[el][0].mass
        name = elements[el][0].name.capitalize()

        this_color = "white"
        if el in census:
            this_color = map_colors[census[el] - 1]
            ndets = r"\textbf{" + str(census[el]) + "}"
            ax.annotate(
                ndets,
                xy=(x + 0.8, y + 0.95),
                xycoords="data",
                size=14,
                color="black",
                ha="right",
                va="top",
            )

        rect = patches.Rectangle(
            (x, y),
            0.9,
            1.05,
            linewidth=1,
            edgecolor="black",
            facecolor=this_color,
            alpha=0.5,
        )
        ax.add_patch(rect)
        ax.annotate(
            num,
            xy=(x + 0.1, y + 0.95),
            xycoords="data",
            size=14,
            color="black",
            ha="left",
            va="top",
        )
        ax.annotate(
            sym,
            xy=(x + 0.1, y + 0.70),
            xycoords="data",
            size=20,
            color="black",
            ha="left",
            va="top",
        )
        ax.annotate(
            mass,
            xy=(x + 0.1, y + 0.42),
            xycoords="data",
            size=8,
            color="black",
            ha="left",
            va="top",
        )
        ax.annotate(
            name,
            xy=(x + 0.1, y + 0.29),
            xycoords="data",
            size=8,
            color="black",
            ha="left",
            va="top",
        )

    plt.axis("equal")
    plt.axis("off")
    plt.show()

    # write out the figure
    plt.savefig(
        filename if filename is not None else "periodic_heatmap.pdf",
        format="pdf",
        transparent=True,
        bbox_inches="tight",
        pad_inches=0,
    )

    # the bit below crops off extra white space.  This only works on Macs with the TexLive pdfcrop utility installed.
    if pdf_crop is True:
        os.system("pdfcrop --margins -0 periodic_heatmap.pdf periodic_heatmap.pdf")


def mass_by_wavelengths(mol_list=None, bw=0.5, filename=None):
    """
    Makes a Kernel Density Estimate plot of detections at each wavelength vs mass.

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)
    bw : float
        The bandwidth to be used for the Kernel Density Estimate (default is 0.5)
    filename : str
        The filename for the output images (default is 'mass_by_wavelengths_kde.pdf')
    """

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    # gather the data
    my_dict = {
        "UV": [],
        "Vis": [],
        "IR": [],
        "sub-mm": [],
        "mm": [],
        "cm": [],
        "UV-Vis": [],
    }

    for x in mol_list:
        for y in my_dict:
            if y in x.wavelengths:
                my_dict[y].append(x.mass)
    for x in my_dict["UV"]:
        if x not in my_dict["UV-Vis"]:
            my_dict["UV-Vis"].append(x)
    for x in my_dict["Vis"]:
        if x not in my_dict["UV-Vis"]:
            my_dict["UV-Vis"].append(x)

    # get the plot set up
    plt.close("Detections at Wavelengths by Mass")
    fig = plt.figure(num="Detections at Wavelengths by Mass", figsize=(10, 8))
    plt.ion()

    # set some font defaults
    fontparams = {"size": 24, "family": "sans-serif", "sans-serif": ["Helvetica"]}
    plt.rc("font", **fontparams)
    plt.rc("mathtext", fontset="stixsans")

    # load up an axis
    ax = fig.add_subplot(111)
    ax.tick_params(axis="x", which="both", direction="in", length=5, width=1)
    ax.tick_params(axis="y", which="both", direction="in", length=5, width=1)

    # label the axes
    plt.xlabel("Atomic Mass (amu)")
    plt.ylabel("Probability Density Estimate")

    # Do the estimates and plot them
    xvals = np.arange(0, 160)

    density_cm = gkde(my_dict["cm"])
    density_cm.covariance_factor = lambda: bw
    density_cm._compute_covariance()

    density_mm = gkde(my_dict["mm"])
    density_mm.covariance_factor = lambda: bw
    density_mm._compute_covariance()

    density_submm = gkde(my_dict["sub-mm"])
    density_submm.covariance_factor = lambda: bw
    density_submm._compute_covariance()

    density_IR = gkde(my_dict["IR"])
    density_IR.covariance_factor = lambda: bw
    density_IR._compute_covariance()

    density_UV = gkde(my_dict["UV-Vis"])
    density_UV.covariance_factor = lambda: bw
    density_UV._compute_covariance()

    ax.plot(xvals, density_cm(xvals), color="dodgerblue")
    ax.fill_between(xvals, density_cm(xvals), 0, facecolor="dodgerblue", alpha=0.25, zorder=4)
    ax.annotate(
        "{}".format(len(my_dict["cm"])),
        xy=(78, 0.01),
        xycoords="data",
        ha="left",
        va="bottom",
        color="dodgerblue",
    )

    max_mm = max([x for x in my_dict["mm"] if x < 160])
    ax.plot(xvals[: max_mm + 1], density_mm(xvals[: max_mm + 1]), color="darkorange")
    ax.fill_between(
        xvals[: max_mm + 1],
        density_mm(xvals[: max_mm + 1]),
        0,
        facecolor="darkorange",
        alpha=0.25,
    )
    ax.annotate(
        "{}".format(len(my_dict["mm"])),
        xy=(54, 0.022),
        xycoords="data",
        ha="left",
        va="bottom",
        color="darkorange",
    )

    ax.plot(xvals, density_submm(xvals), color="forestgreen")
    ax.fill_between(xvals, density_submm(xvals), 0, facecolor="forestgreen", alpha=0.25)
    ax.annotate(
        "{}".format(len(my_dict["sub-mm"])),
        xy=(40, 0.038),
        xycoords="data",
        ha="left",
        va="bottom",
        color="forestgreen",
    )

    max_IR = max([x for x in my_dict["IR"] if x < 160])
    ax.plot(xvals[: max_IR + 1], density_IR(xvals[: max_IR + 1]), color="black")
    ax.fill_between(
        xvals[: max_IR + 1],
        density_IR(xvals[: max_IR + 1]),
        0,
        facecolor="black",
        alpha=0.25,
        zorder=5,
    )
    ax.annotate(
        "{}".format(len(my_dict["IR"])),
        xy=(68, 0.0025),
        xycoords="data",
        ha="left",
        va="bottom",
        color="black",
    )

    max_UV = max(my_dict["UV-Vis"])
    ax.plot(xvals[: max_UV + 1], density_UV(xvals[: max_UV + 1]), color="violet")
    ax.fill_between(
        xvals[: max_UV + 1],
        density_UV(xvals[: max_UV + 1]),
        0,
        facecolor="violet",
        alpha=0.25,
    )
    ax.annotate(
        "{}".format(len(my_dict["UV-Vis"])),
        xy=(14, 0.033),
        xycoords="data",
        ha="left",
        va="bottom",
        color="violet",
    )

    ax.annotate(
        r"\underline{Detection Wavelengths}",
        xy=(0.95, 0.97),
        xycoords="axes fraction",
        color="black",
        ha="right",
        va="top",
    )
    ax.annotate(
        "centimeter",
        xy=(0.95, 0.9),
        xycoords="axes fraction",
        color="dodgerblue",
        ha="right",
        va="top",
    )
    ax.annotate(
        "millimeter",
        xy=(0.95, 0.85),
        xycoords="axes fraction",
        color="darkorange",
        ha="right",
        va="top",
    )
    ax.annotate(
        "sub-millimeter",
        xy=(0.95, 0.8),
        xycoords="axes fraction",
        color="forestgreen",
        ha="right",
        va="top",
    )
    ax.annotate(
        "infrared",
        xy=(0.95, 0.75),
        xycoords="axes fraction",
        color="black",
        ha="right",
        va="top",
    )
    ax.annotate(
        "visible/ultraviolet",
        xy=(0.95, 0.7),
        xycoords="axes fraction",
        color="violet",
        ha="right",
        va="top",
    )

    plt.savefig(
        filename if filename is not None else "mass_by_wavelengths_kde.pdf",
        format="pdf",
        transparent=True,
        bbox_inches="tight",
        pad_inches=0,
    )

    plt.show()

    return


def mols_waves_by_atoms(mol_list=None, bw=0.5, filename=None):
    """
    Makes plots of molecules detected in each wavelength range by number of atoms, excluding
    fullerenes.  For plots with sufficient datapoints, a Kernel Density Estimate is used, otherwise
    a histogram is provided.

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)
    bw : float
        The bandwidth to be used for the Kernel Density Estimate (default is 0.5)
    filename : str
        The filename for the output images (default is 'mols_waves_by_atoms.pdf')
    """

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    # Close an old figure if it exists and initialize a new figure
    plt.close("Molecules Detected in Each Wavelength by Number of Atoms")
    plt.figure(num="Molecules Detected in Each Wavelength by Number of Atoms", figsize=(10, 8))
    plt.ion()

    # set some font defaults
    fontparams = {"size": 18, "family": "sans-serif", "sans-serif": ["Helvetica"]}
    plt.rc("font", **fontparams)
    plt.rc("mathtext", fontset="stixsans")

    # gather the data
    my_dict = {
        "UV": [],
        "Vis": [],
        "IR": [],
        "sub-mm": [],
        "mm": [],
        "cm": [],
    }

    max_n = []
    for x in mol_list:
        for y in my_dict:
            if y in x.wavelengths:
                if x.fullerene is True:
                    continue
                else:
                    my_dict[y].append(x.natoms)
                    max_n.append(x.natoms)

    max_n = max(max_n)

    # pull up some axes to plot on
    ax1 = plt.subplot(231)
    ax2 = plt.subplot(232)
    ax3 = plt.subplot(233)
    ax4 = plt.subplot(234)
    ax5 = plt.subplot(235)
    ax6 = plt.subplot(236)

    xvals = np.arange(0, max_n + 1, 0.5)

    ax1.tick_params(axis="x", which="both", direction="in", length=5, width=1)
    ax1.tick_params(axis="y", which="both", direction="in", length=5, width=1)
    ax2.tick_params(axis="x", which="both", direction="in", length=5, width=1)
    ax2.tick_params(axis="y", which="both", direction="in", length=5, width=1)
    ax3.tick_params(axis="x", which="both", direction="in", length=5, width=1)
    ax3.tick_params(axis="y", which="both", direction="in", length=5, width=1)
    ax4.tick_params(axis="x", which="both", direction="in", length=5, width=1)
    ax4.tick_params(axis="y", which="both", direction="in", length=5, width=1)
    ax5.tick_params(axis="x", which="both", direction="in", length=5, width=1)
    ax5.tick_params(axis="y", which="both", direction="in", length=5, width=1)
    ax6.tick_params(axis="x", which="both", direction="in", length=5, width=1)
    ax6.tick_params(axis="y", which="both", direction="in", length=5, width=1)

    density_cm = gkde(my_dict["cm"])
    density_cm.covariance_factor = lambda: bw
    density_cm._compute_covariance()
    ax1.plot(xvals, density_cm(xvals))
    ax1.fill_between(xvals, density_cm(xvals), 0, facecolor="dodgerblue", alpha=0.25)
    ax1.annotate("cm", xy=(0.95, 0.96), xycoords="axes fraction", ha="right", va="top", size=24)
    ax1.set_ylim([0, 0.65])
    ax1.set_xticks([0, 5, 10, 15, 20])
    ax1.set_xticklabels([])
    ax1.set_ylabel("Probability Density Estimate")

    density_mm = gkde(my_dict["mm"])
    density_mm.covariance_factor = lambda: bw
    density_mm._compute_covariance()
    ax2.plot(xvals, density_mm(xvals))
    ax2.fill_between(xvals, density_mm(xvals), 0, facecolor="dodgerblue", alpha=0.25)
    ax2.annotate("mm", xy=(0.95, 0.96), xycoords="axes fraction", ha="right", va="top", size=24)
    ax2.set_ylim([0, 0.65])
    ax2.set_xticks([0, 5, 10, 15, 20])
    ax2.set_xticklabels([])

    density_submm = gkde(my_dict["sub-mm"])
    density_submm.covariance_factor = lambda: bw
    density_submm._compute_covariance()
    ax3.plot(xvals, density_submm(xvals))
    ax3.fill_between(xvals, density_submm(xvals), 0, facecolor="dodgerblue", alpha=0.25)
    ax3.annotate("sub-mm", xy=(0.95, 0.96), xycoords="axes fraction", ha="right", va="top", size=24)
    ax3.set_ylim([0, 0.65])
    ax3.set_xticks([0, 5, 10, 15, 20])
    ax3.set_xticklabels([])

    density_IR = gkde(my_dict["IR"])
    density_IR.covariance_factor = lambda: bw
    density_IR._compute_covariance()
    ax4.plot(xvals, density_IR(xvals))
    ax4.fill_between(xvals, density_IR(xvals), 0, facecolor="dodgerblue", alpha=0.25)
    ax4.annotate("IR", xy=(0.95, 0.96), xycoords="axes fraction", ha="right", va="top", size=24)
    ax4.set_ylim([0, 0.65])
    ax4.set_xticks([0, 5, 10, 15, 20])
    ax4.set_xlabel(r"\# of Atoms")
    ax4.set_ylabel("Probability Density Estimate")

    ax5.hist(my_dict["Vis"], bins=[0.5, 1.5, 2.5, 3.5, 4.5])
    ax5.annotate("Vis", xy=(0.95, 0.96), xycoords="axes fraction", ha="right", va="top", size=24)
    ax5.set_xlim([0, 20])
    ax5.set_ylim([0, 8])
    ax5.set_xticks([0, 5, 10, 15, 20])
    ax5.set_xticklabels([])
    ax5.set_ylabel(r"\# of Detected Molecules")

    ax6.hist(my_dict["UV"], bins=[0.5, 1.5, 2.5, 3.5, 4.5])
    ax6.annotate("UV", xy=(0.95, 0.96), xycoords="axes fraction", ha="right", va="top", size=24)
    ax6.set_xlim([0, 20])
    ax6.set_ylim([0, 8])
    ax6.set_xticks([0, 5, 10, 15, 20])
    ax6.set_xticklabels([])

    plt.tight_layout()
    plt.show()
    plt.savefig(
        filename if filename is not None else "mols_waves_by_atoms.pdf",
        format="pdf",
        transparent=True,
        bbox_inches="tight",
        pad_inches=0,
    )

    return


def du_histogram(mol_list=None, filename=None):
    """
    Makes a histogram of the degree of unsaturation of molecules for which that value has been
    calculated.

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)
    filename : str
        The filename for the output images (default is 'du_histogram.pdf')
    """

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    # Close an old figure if it exists and initialize a new figure
    plt.close("Degree of Unsaturation Histogram")
    plt.figure(num="Degree of Unsaturation Histogram", figsize=(10, 8))
    plt.ion()

    # set some font defaults
    fontparams = {"size": 24, "family": "sans-serif", "sans-serif": ["Helvetica"]}
    plt.rc("font", **fontparams)
    plt.rc("mathtext", fontset="stixsans")

    # gather the data
    dus = []
    for x in mol_list:
        # make sure we're working with molecules containing only H, D, N, C, Cl, F, O, or S.
        atoms = []
        for atom in ["H", "D", "N", "C", "Cl", "F", "S", "O"]:
            if atom in x.atoms:
                atoms.append(x.atoms[atom])
        if np.sum(atoms) == x.natoms and x.du is not None and x.fullerene is not True:
            dus.append(x.du)

    # set up a plot
    ax = plt.subplot(111)

    ax.tick_params(axis="x", which="both", direction="in", length=5, width=1)
    ax.tick_params(axis="y", which="both", direction="in", length=5, width=1)

    ax.yaxis.set_ticks_position("both")
    ax.xaxis.set_ticks_position("both")

    plt.xlabel("Degree of Unsaturation")
    plt.ylabel(r"\# of Detected Molecules")

    bins = np.arange(-0.25, 12.5, 0.5)
    (n, bins, _) = ax.hist(dus, bins=bins, facecolor="dodgerblue", alpha=0.25)
    ax.hist(dus, bins=bins, edgecolor="royalblue", linewidth=1.5, fill=False)
    ax.set_ylim(0, 35)

    ax.annotate(
        r"\ce{CH4}, \ce{CH3OH}, \ce{CH3Cl}, ...",
        xy=(0, n[0] + 1),
        xycoords="data",
        rotation=90,
        size=16,
        ha="center",
        va="bottom",
    )
    ax.annotate(r"\ce{HC11N}", xy=(12, n[12 * 2] + 1), xycoords="data", rotation=90, size=16, ha="center", va="bottom")

    plt.tight_layout()
    plt.show()

    plt.savefig(
        filename if filename is not None else "du_histogram.pdf", format="pdf", transparent=True, bbox_inches="tight"
    )

    return


def type_pie_chart(mol_list=None, filename=None):
    """
    Makes a pie chart of the fraction of interstellar molecules that are 
    neutral, radical, cation, cyclic, pahs, fullerenes, or anions

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)
    filename : str
        The filename for the output images (default is 'type_pie_chart.pdf')
    """

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    # Close an old figure if it exists and initialize a new figure
    plt.close("Type Pie Chart")
    plt.figure(num="Type Pie Chart", figsize=(10, 8))
    plt.ion()

    # set some font defaults
    fontparams = {"size": 24, "family": "sans-serif", "sans-serif": ["Helvetica"]}
    plt.rc("font", **fontparams)
    plt.rc("mathtext", fontset="stixsans")

    # gather the data
    my_dict = {
        "Neutral": [0],
        "Radical": [0],
        "Cation": [0],
        "Cyclic": [0],
        "Anion": [0],
        "Fullerene": [0],
        "PAH": [0],
    }

    for mol in mol_list:
        if mol.neutral is True:
            my_dict["Neutral"][0] += 1
        if mol.radical is True:
            my_dict["Radical"][0] += 1
        if mol.cation is True:
            my_dict["Cation"][0] += 1
        if mol.cyclic is True:
            my_dict["Cyclic"][0] += 1
        if mol.anion is True:
            my_dict["Anion"][0] += 1
        if mol.fullerene is True:
            my_dict["Fullerene"][0] += 1
        if mol.pah is True:
            my_dict["PAH"][0] += 1

    nmols = len(mol_list)
    for type in my_dict:
        my_dict[type].append(my_dict[type][0] / nmols)

    labels = ["Neutral", "Radical", "Cation", "Cyclic", "Anion", "Fullerene", "PAH"]
    fracs = [my_dict[x][1] for x in labels]

    # set up a plot
    ax = plt.subplot(111)
    size = 0.1

    def getshift(x):
        return -(90 - (360 - 360 * x) / 2)

    ax.pie(
        [fracs[0], 1.0 - fracs[0]],
        colors=["dodgerblue", "#EEEEEE"],
        radius=1,
        startangle=getshift(fracs[0]),
        wedgeprops=dict(width=size, edgecolor="w", linewidth=1),
    )
    ax.pie(
        [fracs[1], 1.0 - fracs[1]],
        colors=["darkorange", "#EEEEEE"],
        radius=1 - size - 0.02,
        startangle=getshift(fracs[1]),
        wedgeprops=dict(width=size, edgecolor="w", linewidth=1),
    )
    ax.pie(
        [fracs[2], 1.0 - fracs[2]],
        colors=["forestgreen", "#EEEEEE"],
        radius=1 - 2 * size - 0.04,
        startangle=getshift(fracs[2]),
        wedgeprops=dict(width=size, edgecolor="w", linewidth=1),
    )
    ax.pie(
        [fracs[3], 1.0 - fracs[3]],
        colors=["violet", "#EEEEEE"],
        radius=1 - 3 * size - 0.06,
        startangle=getshift(fracs[3]),
        wedgeprops=dict(width=size, edgecolor="w", linewidth=1),
    )
    ax.pie(
        [fracs[4], 1.0 - fracs[4]],
        colors=["red", "#EEEEEE"],
        radius=1 - 4 * size - 0.08,
        startangle=getshift(fracs[4]),
        wedgeprops=dict(width=size, edgecolor="w", linewidth=1),
    )
    ax.pie(
        [fracs[5], 1.0 - fracs[5]],
        colors=["goldenrod", "#EEEEEE"],
        radius=1 - 5 * size - 0.1,
        startangle=getshift(fracs[5]),
        wedgeprops=dict(width=size, edgecolor="w", linewidth=1),
    )
    ax.pie(
        [fracs[6], 1.0 - fracs[6]],
        colors=["royalblue", "#EEEEEE"],
        radius=1 - 6 * size - 0.12,
        startangle=getshift(fracs[6]),
        wedgeprops=dict(width=size, edgecolor="w", linewidth=1),
    )

    ax.annotate(
        r"\textbf{Neutrals}", xy=(0.5, 0.11), xycoords="axes fraction", color="dodgerblue", ha="center", size=14
    )
    ax.annotate(
        r"\textbf{Radicals}", xy=(0.5, 0.16), xycoords="axes fraction", color="darkorange", ha="center", size=14
    )
    ax.annotate(
        r"\textbf{Cations}", xy=(0.5, 0.205), xycoords="axes fraction", color="forestgreen", ha="center", size=14
    )
    ax.annotate(r"\textbf{Cyclics}", xy=(0.5, 0.255), xycoords="axes fraction", color="violet", ha="center", size=14)
    ax.annotate(r"\textbf{Anions}", xy=(0.5, 0.305), xycoords="axes fraction", color="red", ha="center", size=14)
    ax.annotate(
        r"\textbf{Fullerenes}", xy=(0.5, 0.3575), xycoords="axes fraction", color="goldenrod", ha="center", size=14
    )
    ax.annotate(r"\textbf{PAHs}", xy=(0.5, 0.40), xycoords="axes fraction", color="royalblue", ha="center", size=14)

    percents = [r"\textbf{" + "{:.1f}".format((x * 100)) + r"}\%" for x in fracs]

    start = 0.585
    shift = 0.0485
    ax.annotate(
        percents[0],
        xy=(start + 6 * shift, 0.5),
        xycoords="axes fraction",
        color="white",
        ha="center",
        va="center",
        size=12,
        rotation=-90,
    )
    ax.annotate(
        percents[1],
        xy=(start + 5 * shift, 0.5),
        xycoords="axes fraction",
        color="darkorange",
        ha="center",
        va="center",
        size=12,
        rotation=-90,
    )
    ax.annotate(
        percents[2],
        xy=(start + 4 * shift, 0.5),
        xycoords="axes fraction",
        color="forestgreen",
        ha="center",
        va="center",
        size=12,
        rotation=-90,
    )
    ax.annotate(
        percents[3],
        xy=(start + 3 * shift, 0.5),
        xycoords="axes fraction",
        color="violet",
        ha="center",
        va="center",
        size=12,
        rotation=-90,
    )
    ax.annotate(
        percents[4],
        xy=(start + 2 * shift, 0.5),
        xycoords="axes fraction",
        color="red",
        ha="center",
        va="center",
        size=12,
        rotation=-90,
    )
    ax.annotate(
        percents[5],
        xy=(start + 1 * shift, 0.5),
        xycoords="axes fraction",
        color="goldenrod",
        ha="center",
        va="center",
        size=12,
        rotation=-90,
    )
    ax.annotate(
        percents[6],
        xy=(start + 0 * shift, 0.5),
        xycoords="axes fraction",
        color="royalblue",
        ha="center",
        va="center",
        size=12,
        rotation=-90,
    )

    plt.tight_layout()
    plt.show()

    plt.savefig(
        filename if filename is not None else "type_pie_chart.pdf", format="pdf", transparent=True, bbox_inches="tight"
    )  # ,pad_inches=-.65)

    return


def source_pie_chart(mol_list=None, filename=None, format='pdf'):
    """
    Makes a pie chart of the fraction of interstellar molecules detected 
    in carbon stars, dark clouds, los clouds, star forming regions, 
    and other types of sources

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)
    filename : str
        The filename for the output images (default is 'source_pie_chart.pdf')
    """ 

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    # Close an old figure if it exists and initialize a new figure
    plt.close("Source Pie Chart")
    plt.figure(num="Source Pie Chart", figsize=(10, 8))
    plt.ion()

    # set some font defaults
    fontparams = {"size": 24, "family": "sans-serif", "sans-serif": ["Helvetica"]}
    plt.rc("font", **fontparams)
    plt.rc("mathtext", fontset="stixsans")

    # gather the data
    my_dict = {
        "Carbon Star": [0],
        "Dark Cloud": [0],
        "LOS Cloud": [0],
        "SFR": [0],
        "Other": [0],
    }

    # we have to be a little careful here, because for a given source, there can be two SFRs listed, and we only want to credit it once
    for mol in mol_list:

        # we'll make a dictionary here to flag if we've credited things already
        credit_dict = {
            "Carbon Star": False,
            "Dark Cloud": False,
            "LOS Cloud": False,
            "SFR": False,
            "Other": False,
        }

        for source in mol.sources:
            if source.type in my_dict:
                if credit_dict[source.type] is False:
                    my_dict[source.type][0] += 1
                    credit_dict[source.type] = True
            else:
                if credit_dict["Other"] is False:
                    my_dict["Other"][0] += 1
                    credit_dict["Other"] = True

    nmols = len(mol_list)
    for type in my_dict:
        my_dict[type].append(my_dict[type][0] / nmols)

    labels = ["SFR", "Carbon Star", "Dark Cloud", "Other", "LOS Cloud"]
    fracs = [my_dict[x][1] for x in labels]

    # set up a plot
    ax = plt.subplot(111)
    size = 0.1

    def getshift(x):
        return -(90 - (360 - 360 * x) / 2)

    ax.pie(
        [fracs[0], 1.0 - fracs[0]],
        colors=["dodgerblue", "#EEEEEE"],
        radius=1 - 0 * size - 0 * 0.02,
        startangle=getshift(fracs[0]),
        wedgeprops=dict(width=size, edgecolor="w", linewidth=1),
    )
    ax.pie(
        [fracs[2], 1.0 - fracs[2]],
        colors=["forestgreen", "#EEEEEE"],
        radius=1 - 1 * size - 1 * 0.02,
        startangle=getshift(fracs[2]),
        wedgeprops=dict(width=size, edgecolor="w", linewidth=1),
    )
    ax.pie(
        [fracs[1], 1.0 - fracs[1]],
        colors=["darkorange", "#EEEEEE"],
        radius=1 - 2 * size - 2 * 0.02,
        startangle=getshift(fracs[1]),
        wedgeprops=dict(width=size, edgecolor="w", linewidth=1),
    )
    ax.pie(
        [fracs[3], 1.0 - fracs[3]],
        colors=["violet", "#EEEEEE"],
        radius=1 - 3 * size - 3 * 0.02,
        startangle=getshift(fracs[3]),
        wedgeprops=dict(width=size, edgecolor="w", linewidth=1),
    )
    ax.pie(
        [fracs[4], 1.0 - fracs[4]],
        colors=["red", "#EEEEEE"],
        radius=1 - 4 * size - 4 * 0.02,
        startangle=getshift(fracs[4]),
        wedgeprops=dict(width=size, edgecolor="w", linewidth=1),
    )

    ax.annotate(r"\textbf{SFR}", xy=(0.5, 0.11), xycoords="axes fraction", color="dodgerblue", ha="center", size=14)
    ax.annotate(
        r"\textbf{Carbon Star}", xy=(0.5, 0.205), xycoords="axes fraction", color="darkorange", ha="center", size=14
    )
    ax.annotate(
        r"\textbf{Dark Cloud}", xy=(0.5, 0.16), xycoords="axes fraction", color="forestgreen", ha="center", size=14
    )
    ax.annotate(r"\textbf{Other}", xy=(0.5, 0.255), xycoords="axes fraction", color="violet", ha="center", size=14)
    ax.annotate(r"\textbf{LOS Cloud}", xy=(0.5, 0.305), xycoords="axes fraction", color="red", ha="center", size=14)

    percents = [r"\textbf{" + "{:.1f}".format((x * 100)) + r"}\%" for x in fracs]

    # Percents indexing:
    # 0 : SFR
    # 1 : Carbon Star
    # 2 : Dark Cloud
    # 3 : Other
    # 4 : LOS Cloud

    ax.annotate(
        percents[4],
        xy=(0.5, 0.68),
        xycoords="axes fraction",
        color="white",
        ha="center",
        size=12,
    )
    ax.annotate(
        percents[3],
        xy=(0.5, 0.725),
        xycoords="axes fraction",
        color="white",
        ha="center",
        size=12,
    )
    ax.annotate(
        percents[1],
        xy=(0.5, 0.775),
        xycoords="axes fraction",
        color="white",
        ha="center",
        size=12,
    )
    ax.annotate(
        percents[2],
        xy=(0.5, 0.825),
        xycoords="axes fraction",
        color="white",
        ha="center",
        size=12,
    )
    ax.annotate(
        percents[0],
        xy=(0.5, 0.87),
        xycoords="axes fraction",
        color="white",
        ha="center",
        size=12,
    )

    plt.tight_layout()

    if format=='pdf':

        plt.savefig(
            filename if filename is not None else "source_pie_chart.pdf",
            format="pdf",
            transparent=True,
            bbox_inches="tight",
            pad_inches=-0.65,
        )

    elif format=='png':
    
        plt.savefig(
            filename if filename is not None else "source_pie_chart.png",
            format="png",
            dpi=600,
            transparent=True,
            bbox_inches="tight",
            pad_inches=-0.65,
        )   

    plt.show()

    return


def indiv_source_pie_chart(mol_list=None, filename=None, format='pdf'):
    """
    Makes a pie chart of the fraction of interstellar molecules 
    detected in IRC+10216, TMC-1, Orion, and Sgr

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)
    filename : str
        The filename for the output images (default is 'indiv_source_pie_chart.pdf')
    """ 

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    # Close an old figure if it exists and initialize a new figure
    plt.close("Individual Source Pie Chart")
    plt.figure(num="Individual Source Pie Chart", figsize=(10, 8))
    plt.ion()

    # set some font defaults
    fontparams = {"size": 24, "family": "sans-serif", "sans-serif": ["Helvetica"]}
    plt.rc("font", **fontparams)
    plt.rc("mathtext", fontset="stixsans")

    # gather the data
    my_dict = {
        "IRC+10216": [0],
        "TMC-1": [0],
        "Orion": [0],
        "Sgr B2": [0],
        "Other": [0],
    }

    for mol in mol_list:
        for source in mol.sources:
            other = False
            if source == IRC10216:
                my_dict["IRC+10216"][0] += 1
                other = True
            if source == TMC1:
                my_dict["TMC-1"][0] += 1
                other = True
            if source == Orion:
                my_dict["Orion"][0] += 1
                other = True
            if source == SgrB2:
                my_dict["Sgr B2"][0] += 1
                other = True
            if other is False:
                my_dict["Other"][0] += 1

    nmols = len(mol_list)
    for type in my_dict:
        my_dict[type].append(my_dict[type][0] / nmols)
    labels = ["Other", "Sgr B2", "IRC+10216", "TMC-1", "Orion"]
    fracs = [my_dict[x][1] for x in labels]

    # set up a plot
    ax = plt.subplot(111)
    size = 0.1

    def getshift(x):
        return -(90 - (360 - 360 * x) / 2)

    ax.pie(
        [fracs[0], 1.0 - fracs[0]],
        colors=["dodgerblue", "#EEEEEE"],
        radius=1 - 0 * size - 0 * 0.02,
        startangle=getshift(fracs[0]),
        wedgeprops=dict(width=size, edgecolor="w", linewidth=1),
    )
    ax.pie(
        [fracs[1], 1.0 - fracs[1]],
        colors=["darkorange", "#EEEEEE"],
        radius=1 - 1 * size - 1 * 0.02,
        startangle=getshift(fracs[1]),
        wedgeprops=dict(width=size, edgecolor="w", linewidth=1),
    )
    ax.pie(
        [fracs[2], 1.0 - fracs[2]],
        colors=["forestgreen", "#EEEEEE"],
        radius=1 - 3 * size - 3 * 0.02,
        startangle=getshift(fracs[2]),
        wedgeprops=dict(width=size, edgecolor="w", linewidth=1),
    )
    ax.pie(
        [fracs[3], 1.0 - fracs[3]],
        colors=["violet", "#EEEEEE"],
        radius=1 - 2 * size - 2 * 0.02,
        startangle=getshift(fracs[3]),
        wedgeprops=dict(width=size, edgecolor="w", linewidth=1),
    )
    ax.pie(
        [fracs[4], 1.0 - fracs[4]],
        colors=["red", "#EEEEEE"],
        radius=1 - 4 * size - 4 * 0.02,
        startangle=getshift(fracs[4]),
        wedgeprops=dict(width=size, edgecolor="w", linewidth=1),
    )

    ax.annotate(r"\textbf{Other}", xy=(0.5, 0.11), xycoords="axes fraction", color="dodgerblue", ha="center", size=14)
    ax.annotate(r"\textbf{Sgr B2}", xy=(0.5, 0.16), xycoords="axes fraction", color="darkorange", ha="center", size=14)
    ax.annotate(
        r"\textbf{IRC+10216}", xy=(0.5, 0.255), xycoords="axes fraction", color="forestgreen", ha="center", size=14
    )
    ax.annotate(r"\textbf{TMC-1}", xy=(0.5, 0.205), xycoords="axes fraction", color="violet", ha="center", size=14)
    ax.annotate(r"\textbf{Orion}", xy=(0.5, 0.305), xycoords="axes fraction", color="red", ha="center", size=14)

    percents = [r"\textbf{" + "{:.1f}".format((x * 100)) + r"}\%" for x in fracs]

    # percents indexing
    # 0 : Other
    # 1 : Sgr B2
    # 2 : IRC+10216
    # 3 : TMC-1
    # 4 : Orion

    ax.annotate(
        percents[4],
        xy=(0.51, 0.68),
        xycoords="axes fraction",
        color="white",
        ha="center",
        size=12,
    )
    ax.annotate(
        percents[2],
        xy=(0.51, 0.725),
        xycoords="axes fraction",
        color="white",
        ha="center",
        size=12,
    )
    ax.annotate(
        percents[3],
        xy=(0.51, 0.775),
        xycoords="axes fraction",
        color="white",
        ha="center",
        size=12,
    )
    ax.annotate(
        percents[1],
        xy=(0.51, 0.825),
        xycoords="axes fraction",
        color="white",
        ha="center",
        size=12,
    )
    ax.annotate(
        percents[0],
        xy=(0.51, 0.87),
        xycoords="axes fraction",
        color="white",
        ha="center",
        size=12,
    )

    plt.tight_layout()
    plt.show()

    if format=='pdf':

        plt.savefig(
            filename if filename is not None else "indiv_source_pie_chart.pdf",
            format="pdf",
            transparent=True,
            bbox_inches="tight",
            pad_inches=-0.65,
        )
    
    elif format=='png':
    
        plt.savefig(
            filename if filename is not None else "indiv_source_pie_chart.png",
            format="png",
            dpi=600,
            transparent=True,
            bbox_inches="tight",
            pad_inches=-0.65,
        )   
    

    return


def mol_type_by_source_type(mol_list=None, filename=None):
    """
    Generates four pie charts, one for each generalized source type, 
    with the wedges for the types of molecules detected first in each type

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)
    filename : str
        The filename for the output images (default is 'mol_type_by_source_type.pdf')
    """


    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    # Close an old figure if it exists and initialize a new figure
    plt.close("Molecule Type by Source Type")
    fig, axs = plt.subplots(2, 2, num="Molecule Type by Source Type", figsize=(15, 12))
    plt.ion()

    # set some font defaults
    fontparams = {"size": 26, "family": "sans-serif", "sans-serif": ["Helvetica"]}
    plt.rc("font", **fontparams)
    plt.rc("mathtext", fontset="stixsans")

    # collect the data
    type_dict = {
        "Carbon Star": {"Anion": 0, "Cation": 0, "Cyclic": 0, "Neutral": 0, "Radical": 0},
        "Dark Cloud": {"Anion": 0, "Cation": 0, "Cyclic": 0, "Neutral": 0, "Radical": 0},
        "LOS Cloud": {"Anion": 0, "Cation": 0, "Cyclic": 0, "Neutral": 0, "Radical": 0},
        "SFR": {"Anion": 0, "Cation": 0, "Cyclic": 0, "Neutral": 0, "Radical": 0},
    }

    for mol in mol_list:
        # we'll make a dictionary here to flag if we've credited things already
        credit_dict = {
            "Carbon Star": False,
            "Dark Cloud": False,
            "LOS Cloud": False,
            "SFR": False,
        }

        for source in mol.sources:
            if source.type in type_dict:
                if credit_dict[source.type] is False:
                    if mol.anion is True:
                        type_dict[source.type]["Anion"] += 1
                    if mol.cation is True:
                        type_dict[source.type]["Cation"] += 1
                    if mol.cyclic is True:
                        type_dict[source.type]["Cyclic"] += 1
                    if mol.neutral is True:
                        type_dict[source.type]["Neutral"] += 1
                    if mol.radical is True:
                        type_dict[source.type]["Radical"] += 1
                    credit_dict[source.type] = True

    # make the pie charts
    # Carbon Stars
    carbon_data = [
        type_dict["Carbon Star"]["Anion"],
        # type_dict['Carbon Star']['Cation'],
        type_dict["Carbon Star"]["Cyclic"],
        type_dict["Carbon Star"]["Neutral"],
        type_dict["Carbon Star"]["Radical"],
    ]

    carbon_colors = [
        "darkorange",
        #'forestgreen',
        "violet",
        "dodgerblue",
        "red",
    ]

    # Dark Clouds
    dark_data = [
        type_dict["Dark Cloud"]["Anion"],
        type_dict["Dark Cloud"]["Cation"],
        type_dict["Dark Cloud"]["Cyclic"],
        type_dict["Dark Cloud"]["Neutral"],
        type_dict["Dark Cloud"]["Radical"],
    ]

    dark_colors = [
        "darkorange",
        "forestgreen",
        "violet",
        "dodgerblue",
        "red",
    ]

    # LOS Clouds
    los_data = [
        # type_dict['LOS Cloud']['Anion'],
        type_dict["LOS Cloud"]["Cation"],
        # type_dict['LOS Cloud']['Cyclic'],
        type_dict["LOS Cloud"]["Neutral"],
        type_dict["LOS Cloud"]["Radical"],
    ]

    los_colors = [
        #'darkorange',
        "forestgreen",
        #'violet',
        "dodgerblue",
        "red",
    ]

    # SFRs
    sfr_data = [
        # type_dict['SFR']['Anion'],
        type_dict["SFR"]["Cation"],
        type_dict["SFR"]["Cyclic"],
        type_dict["SFR"]["Neutral"],
        type_dict["SFR"]["Radical"],
    ]

    sfr_colors = [
        #'darkorange',
        "forestgreen",
        "violet",
        "dodgerblue",
        "red",
    ]

    types = [
        "Anion",
        "Cation",
        "Cyclic",
        "Neutral",
        "Radical",
    ]

    axs[0, 0].pie(
        carbon_data,
        labels=carbon_data,
        colors=carbon_colors,
        labeldistance=0.8,
        wedgeprops={"linewidth": 1.0, "edgecolor": "black", "alpha": 0.5},
    )
    axs[0, 0].annotate(
        r"\textbf{Carbon Stars}", xy=(0.5, 1.0), xycoords="axes fraction", color="black", ha="center", va="top", size=30
    )

    wedges, _ = axs[0, 1].pie(
        dark_data,
        labels=dark_data,
        colors=dark_colors,
        labeldistance=0.8,
        wedgeprops={"linewidth": 1.0, "edgecolor": "black", "alpha": 0.5},
    )
    axs[0, 1].annotate(
        r"\textbf{Dark Clouds}", xy=(0.5, 1.0), xycoords="axes fraction", color="black", ha="center", va="top", size=30
    )

    axs[1, 0].pie(
        los_data,
        labels=los_data,
        colors=los_colors,
        labeldistance=0.8,
        wedgeprops={"linewidth": 1.0, "edgecolor": "black", "alpha": 0.5},
    )
    axs[1, 0].annotate(
        r"\textbf{LOS Clouds}", xy=(0.5, 1.0), xycoords="axes fraction", color="black", ha="center", va="top", size=30
    )

    axs[1, 1].pie(
        sfr_data,
        labels=sfr_data,
        colors=sfr_colors,
        labeldistance=0.8,
        wedgeprops={"linewidth": 1.0, "edgecolor": "black", "alpha": 0.5},
    )
    axs[1, 1].annotate(
        r"\textbf{SFRs}", xy=(0.5, 1.0), xycoords="axes fraction", color="black", ha="center", va="top", size=30
    )

    axs[0, 1].legend(wedges, types, title="Molecule Types", loc="center left", bbox_to_anchor=(1, 0, 0.5, 1))

    fig.tight_layout()
    fig.subplots_adjust(wspace=-0.3, hspace=0)
    plt.show()

    plt.savefig(
        filename if filename is not None else "mol_type_by_source_type.pdf",
        format="pdf",
        transparent=True,
        bbox_inches="tight",
    )

    return


def du_by_source_type(mol_list=None, bw=0.5, filename=None):
    """
    Makes a Kernel Density Estimate plot of the degrees of unsaturation in each source type

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)
    bw : float
        The bandwidth to be used for the Kernel Density Estimate (default is 0.5)        
    filename : str
        The filename for the output images (default is 'du_by_source_type_kde.pdf')
    """

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    # gather the data
    my_dict = {
        "Carbon Star": [],
        "Dark Cloud": [],
        "LOS Cloud": [],
        "SFR": [],
    }

    # we have to be a little careful here, because for a given source, there can be two SFRs listed, and we only want to credit it once
    for mol in mol_list:
        # we'll make a dictionary here to flag if we've credited things already
        credit_dict = {
            "Carbon Star": False,
            "Dark Cloud": False,
            "LOS Cloud": False,
            "SFR": False,
        }

        if mol.du is None or mol.fullerene is True:
            continue
        for source in mol.sources:
            if source.type in my_dict:
                if credit_dict[source.type] is False:
                    my_dict[source.type].append(mol.du)
                    credit_dict[source.type] = True

    plt.close("DU by Source Type")
    fig = plt.figure(num="DU by Source Type", figsize=(10, 8))
    plt.ion()

    # set some font defaults
    fontparams = {"size": 24, "family": "sans-serif", "sans-serif": ["Helvetica"]}
    plt.rc("font", **fontparams)
    plt.rc("mathtext", fontset="stixsans")

    # load up an axis
    ax = fig.add_subplot(111)
    ax.tick_params(axis="x", which="both", direction="in", length=5, width=1)
    ax.tick_params(axis="y", which="both", direction="in", length=5, width=1)

    # label the axes
    plt.xlabel("Degree of Unsaturation")
    plt.ylabel("Probably Density Estimate")

    # Do the estimates and plot them
    xvals = np.arange(0, 15, 0.1)

    density_carbon = gkde(my_dict["Carbon Star"])
    density_carbon.covariance_factor = lambda: bw
    density_carbon._compute_covariance()

    density_dark = gkde(my_dict["Dark Cloud"])
    density_dark.covariance_factor = lambda: bw
    density_dark._compute_covariance()

    density_los = gkde(my_dict["LOS Cloud"])
    density_los.covariance_factor = lambda: bw
    density_los._compute_covariance()

    density_sfr = gkde(my_dict["SFR"])
    density_sfr.covariance_factor = lambda: bw
    density_sfr._compute_covariance()

    x_ann = 0.97
    y_ann = 0.96
    y_sep = 0.06

    ax.plot(xvals, density_carbon(xvals), color="darkorange")
    ax.fill_between(xvals, density_carbon(xvals), 0, facecolor="darkorange", alpha=0.25, zorder=4)

    ax.plot(xvals, density_dark(xvals), color="forestgreen")
    ax.fill_between(xvals, density_dark(xvals), 0, facecolor="forestgreen", alpha=0.25, zorder=4)

    ax.plot(xvals, density_los(xvals), color="red")
    ax.fill_between(xvals, density_los(xvals), 0, facecolor="red", alpha=0.25, zorder=4)

    ax.plot(xvals, density_sfr(xvals), color="dodgerblue")
    ax.fill_between(xvals, density_sfr(xvals), 0, facecolor="dodgerblue", alpha=0.25, zorder=4)

    ax.annotate(
        r"\underline{Source Types}",
        xy=(x_ann, y_ann - 0 * y_sep),
        xycoords="axes fraction",
        color="black",
        ha="right",
        va="top",
    )

    ax.annotate("SFR", xy=(x_ann, y_ann - y_sep), xycoords="axes fraction", color="dodgerblue", ha="right", va="top")
    ax.annotate(
        "{}".format(len(my_dict["SFR"])), xy=(2.3, 0.3), xycoords="data", ha="left", va="bottom", color="dodgerblue"
    )

    ax.annotate(
        "Carbon Star", xy=(x_ann, y_ann - 2 * y_sep), xycoords="axes fraction", color="darkorange", ha="right", va="top"
    )
    ax.annotate(
        "{}".format(len(my_dict["Carbon Star"])),
        xy=(6, 0.14),
        xycoords="data",
        ha="left",
        va="bottom",
        color="darkorange",
    )

    ax.annotate(
        "Dark Cloud", xy=(x_ann, y_ann - 3 * y_sep), xycoords="axes fraction", color="forestgreen", ha="right", va="top"
    )
    ax.annotate(
        "{}".format(len(my_dict["Dark Cloud"])),
        xy=(10.9, 0.02),
        xycoords="data",
        ha="left",
        va="bottom",
        color="forestgreen",
    )

    ax.annotate("LOS Cloud", xy=(x_ann, y_ann - 4 * y_sep), xycoords="axes fraction", color="red", ha="right", va="top")
    ax.annotate(
        "{}".format(len(my_dict["LOS Cloud"])), xy=(0.65, 0.36), xycoords="data", ha="left", va="bottom", color="red"
    )

    plt.show()

    plt.savefig(
        filename if filename is not None else "du_by_source_type_kde.pdf",
        format="pdf",
        transparent=True,
        bbox_inches="tight",
        pad_inches=0,
    )

    return


def rel_du_by_source_type(mol_list=None, bw=0.5, filename=None):
    """
    Makes a Kernel Density Estimate plot of the relative degrees of 
    unsaturation in each source type

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)
    bw : float
        The bandwidth to be used for the Kernel Density Estimate (default is 0.5)        
    filename : str
        The filename for the output images (default is 'relative_du_by_source_type_kde.pdf')
    """

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    # gather the data
    my_dict = {
        "Carbon Star": [],
        "Dark Cloud": [],
        "LOS Cloud": [],
        "SFR": [],
    }

    # we have to be a little careful here, because for a given source, there can be two SFRs listed, and we only want to credit it once
    for mol in mol_list:
        # we'll make a dictionary here to flag if we've credited things already
        credit_dict = {
            "Carbon Star": False,
            "Dark Cloud": False,
            "LOS Cloud": False,
            "SFR": False,
        }

        if mol.du is None or mol.fullerene is True:
            continue
        for source in mol.sources:
            if source.type in my_dict:
                if credit_dict[source.type] is False:
                    my_dict[source.type].append(mol.du / mol.maxdu)
                    credit_dict[source.type] = True

    plt.close("Relative DU by Source Type")
    _, axs = plt.subplots(2, 2, num="Relative DU by Source Type", figsize=(10, 8))
    plt.ion()

    # set some font defaults
    fontparams = {"size": 18, "family": "sans-serif", "sans-serif": ["Helvetica"]}
    plt.rc("font", **fontparams)
    plt.rc("mathtext", fontset="stixsans")

    # axes ticks and limits
    axs[0, 0].xaxis.set_visible(False)
    axs[0, 0].yaxis.set_visible(False)
    axs[0, 1].xaxis.set_visible(False)
    axs[0, 1].yaxis.set_visible(False)
    axs[1, 1].xaxis.set_visible(False)
    axs[1, 1].yaxis.set_visible(False)

    # label the axes
    axs[1, 0].tick_params(axis="y", which="both", direction="in", length=5, width=1, labelsize=18)
    axs[1, 0].tick_params(axis="x", which="both", direction="in", length=5, width=1, labelsize=18)
    axs[1, 0].set_xlabel("Relative Degree of Unsaturation", fontsize=18)
    axs[1, 0].set_ylabel("Probability Density Estimate", fontsize=18)

    # Do the estimates and plot them
    xvals = np.arange(0, 1, 0.01)

    density_carbon = gkde(my_dict["Carbon Star"])
    density_carbon.covariance_factor = lambda: bw
    density_carbon._compute_covariance()

    density_dark = gkde(my_dict["Dark Cloud"])
    density_dark.covariance_factor = lambda: bw
    density_dark._compute_covariance()

    density_los = gkde(my_dict["LOS Cloud"])
    density_los.covariance_factor = lambda: bw
    density_los._compute_covariance()

    density_sfr = gkde(my_dict["SFR"])
    density_sfr.covariance_factor = lambda: bw
    density_sfr._compute_covariance()

    axs[0, 0].plot(xvals, density_carbon(xvals), color="darkorange")
    axs[0, 0].fill_between(xvals, density_carbon(xvals), 0, facecolor="darkorange", alpha=0.25, zorder=4)
    axs[0, 0].annotate(
        "Carbon Star", xy=[0.04, 0.96], xycoords="axes fraction", ha="left", va="top", size=24, color="darkorange"
    )

    axs[1, 0].plot(xvals, density_dark(xvals), color="forestgreen")
    axs[1, 0].fill_between(xvals, density_dark(xvals), 0, facecolor="forestgreen", alpha=0.25, zorder=4)
    axs[1, 0].annotate(
        "Dark Cloud", xy=[0.04, 0.96], xycoords="axes fraction", ha="left", va="top", size=24, color="forestgreen"
    )

    axs[0, 1].plot(xvals, density_los(xvals), color="red")
    axs[0, 1].fill_between(xvals, density_los(xvals), 0, facecolor="red", alpha=0.25, zorder=4)
    axs[0, 1].annotate(
        "LOS Cloud", xy=[0.04, 0.96], xycoords="axes fraction", ha="left", va="top", size=24, color="red"
    )

    axs[1, 1].plot(xvals, density_sfr(xvals), color="dodgerblue")
    axs[1, 1].fill_between(xvals, density_sfr(xvals), 0, facecolor="dodgerblue", alpha=0.25, zorder=4)
    axs[1, 1].annotate(
        "SFR", xy=[0.04, 0.96], xycoords="axes fraction", ha="left", va="top", size=24, color="dodgerblue"
    )

    plt.subplots_adjust(wspace=0, hspace=0)
    plt.show()

    plt.savefig(
        filename if filename is not None else "relative_du_by_source_type_kde.pdf",
        format="pdf",
        transparent=True,
        bbox_inches="tight",
        pad_inches=0,
    )

    return


def mass_by_source_type(mol_list=None, bw=0.5, filename=None):
    """
    Makes a Kernel Density Estimate plot of the masses in each source type

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)
    bw : float
        The bandwidth to be used for the Kernel Density Estimate (default is 0.5)        
    filename : str
        The filename for the output images (default is 'mass_by_source_type_kde.pdf')
    """

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    # gather the data
    my_dict = {
        "Carbon Star": [],
        "Dark Cloud": [],
        "LOS Cloud": [],
        "SFR": [],
    }

    # we have to be a little careful here, because for a given source, there can be two SFRs listed, and we only want to credit it once
    masses = []
    for mol in mol_list:
        # we'll make a dictionary here to flag if we've credited things already
        credit_dict = {
            "Carbon Star": False,
            "Dark Cloud": False,
            "LOS Cloud": False,
            "SFR": False,
        }

        # drop the fullerenes
        if mol.fullerene is True:
            continue
        # add the mass to the list for axis limit purposes
        masses.append(mol.mass)
        for source in mol.sources:
            if source.type in my_dict:
                if credit_dict[source.type] is False:
                    my_dict[source.type].append(mol.mass)
                    credit_dict[source.type] = True

    plt.close("Mass by Source Type")
    plt.figure(num="Mass by Source Type", figsize=(10, 8))
    plt.ion()

    # set some font defaults
    fontparams = {"size": 24, "family": "sans-serif", "sans-serif": ["Helvetica"]}
    plt.rc("font", **fontparams)
    plt.rc("mathtext", fontset="stixsans")

    # axes ticks and limits
    ax = plt.subplot(111)
    ax.set_xlim([min(masses), max(masses)])
    ax.tick_params(axis="y", which="both", direction="in", length=5, width=1, labelleft="off")
    ax.tick_params(axis="x", which="both", direction="in", length=5, width=1, labelbottom="off")

    # label the axes
    ax.set(xlabel="Molecular Mass (amu)")
    ax.set(ylabel="Probability Density Estimate")

    # Do the estimates and plot them
    xvals = np.arange(0, max(masses), 1)

    density_carbon = gkde(my_dict["Carbon Star"])
    density_carbon.covariance_factor = lambda: 0.5
    density_carbon._compute_covariance()

    density_dark = gkde(my_dict["Dark Cloud"])
    density_dark.covariance_factor = lambda: 0.5
    density_dark._compute_covariance()

    density_los = gkde(my_dict["LOS Cloud"])
    density_los.covariance_factor = lambda: 0.5
    density_los._compute_covariance()

    density_sfr = gkde(my_dict["SFR"])
    density_sfr.covariance_factor = lambda: 0.5
    density_sfr._compute_covariance()

    ax.plot(xvals[min(masses) : max(masses)], density_carbon(xvals[min(masses) : max(masses)]), color="darkorange")
    ax.fill_between(
        xvals[min(masses) : max(masses)],
        density_carbon(xvals[min(masses) : max(masses)]),
        0,
        facecolor="darkorange",
        alpha=0.25,
        zorder=4,
    )

    ax.plot(xvals[min(masses) : max(masses)], density_dark(xvals[min(masses) : max(masses)]), color="forestgreen")
    ax.fill_between(
        xvals[min(masses) : max(masses)],
        density_dark(xvals[min(masses) : max(masses)]),
        0,
        facecolor="forestgreen",
        alpha=0.25,
        zorder=4,
    )

    ax.plot(xvals[min(masses) : max(masses)], density_los(xvals[min(masses) : max(masses)]), color="red")
    ax.fill_between(
        xvals[min(masses) : max(masses)],
        density_los(xvals[min(masses) : max(masses)]),
        0,
        facecolor="red",
        alpha=0.25,
        zorder=4,
    )

    ax.plot(xvals[min(masses) : max(masses)], density_sfr(xvals[min(masses) : max(masses)]), color="dodgerblue")
    ax.fill_between(
        xvals[min(masses) : max(masses)],
        density_sfr(xvals[min(masses) : max(masses)]),
        0,
        facecolor="dodgerblue",
        alpha=0.25,
        zorder=4,
    )

    x_ann = 0.97
    y_ann = 0.96
    y_sep = 0.06

    ax.annotate(
        r"\underline{Source Types}",
        xy=(x_ann, y_ann - 0 * y_sep),
        xycoords="axes fraction",
        color="black",
        ha="right",
        va="top",
    )
    ax.annotate("SFR", xy=(x_ann, y_ann - y_sep), xycoords="axes fraction", color="dodgerblue", ha="right", va="top")
    ax.annotate(
        "Carbon Star", xy=(x_ann, y_ann - 2 * y_sep), xycoords="axes fraction", color="darkorange", ha="right", va="top"
    )
    ax.annotate(
        "Dark Cloud", xy=(x_ann, y_ann - 3 * y_sep), xycoords="axes fraction", color="forestgreen", ha="right", va="top"
    )
    ax.annotate("LOS Cloud", xy=(x_ann, y_ann - 4 * y_sep), xycoords="axes fraction", color="red", ha="right", va="top")

    plt.savefig(
        filename if filename is not None else "mass_by_source_type_kde.pdf",
        format="pdf",
        transparent=True,
        bbox_inches="tight",
        pad_inches=0,
    )

    plt.show()

    return


def waves_by_source_type(mol_list=None, filename=None):
    """
    Generates four pie charts, one for each generalized source type, 
    with the wedges for the wavelengths used for first detections in those sources

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)     
    filename : str
        The filename for the output images (default is 'waves_by_source_type.pdf')
    """

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    # Close an old figure if it exists and initialize a new figure
    plt.close("Wavelength by Source Type")
    fig, axs = plt.subplots(2, 2, num="Wavelength by Source Type", figsize=(15, 12))
    plt.ion()

    # set some font defaults
    fontparams = {"size": 26, "family": "sans-serif", "sans-serif": ["Helvetica"]}
    plt.rc("font", **fontparams)
    plt.rc("mathtext", fontset="stixsans")

    # collect the data
    type_dict = {
        "Carbon Star": {"cm": 0, "mm": 0, "sub-mm": 0, "IR": 0, "UV": 0, "Vis": 0},
        "Dark Cloud": {"cm": 0, "mm": 0, "sub-mm": 0, "IR": 0, "UV": 0, "Vis": 0},
        "LOS Cloud": {"cm": 0, "mm": 0, "sub-mm": 0, "IR": 0, "UV": 0, "Vis": 0},
        "SFR": {"cm": 0, "mm": 0, "sub-mm": 0, "IR": 0, "UV": 0, "Vis": 0},
    }

    for mol in mol_list:
        # we'll make a dictionary here to flag if we've credited things already
        credit_dict = {
            "Carbon Star": False,
            "Dark Cloud": False,
            "LOS Cloud": False,
            "SFR": False,
        }

        for source in mol.sources:
            if source.type in type_dict:
                if credit_dict[source.type] is False:
                    for wave in mol.wavelengths:
                        type_dict[source.type][wave] += 1
                    credit_dict[source.type] = True

    # make the pie charts
    # Carbon Stars
    carbon_data = [
        type_dict["Carbon Star"]["cm"],
        type_dict["Carbon Star"]["mm"],
        type_dict["Carbon Star"]["sub-mm"],
        type_dict["Carbon Star"]["IR"],
        # type_dict['Carbon Star']['UV'] + type_dict['Carbon Star']['Vis'],
    ]

    carbon_colors = [
        "dodgerblue",
        "darkorange",
        "forestgreen",
        "black",
        #'violet',
    ]

    # Dark Clouds
    dark_data = [
        type_dict["Dark Cloud"]["cm"],
        type_dict["Dark Cloud"]["mm"],
        # type_dict['Dark Cloud']['sub-mm'],
        # type_dict['Dark Cloud']['IR'],
        # type_dict['Dark Cloud']['UV'] + type_dict['Dark Cloud']['Vis'],
    ]

    dark_colors = [
        "dodgerblue",
        "darkorange",
        #'forestgreen',
        #'black',
        #'violet',
    ]

    # LOS Clouds

    los_data = [
        type_dict["LOS Cloud"]["cm"],
        type_dict["LOS Cloud"]["mm"],
        type_dict["LOS Cloud"]["sub-mm"],
        type_dict["LOS Cloud"]["IR"],
        type_dict["LOS Cloud"]["UV"] + type_dict["LOS Cloud"]["Vis"],
    ]

    los_colors = [
        "dodgerblue",
        "darkorange",
        "forestgreen",
        "black",
        "violet",
    ]

    # SFRs

    sfr_data = [
        type_dict["SFR"]["cm"],
        type_dict["SFR"]["mm"],
        type_dict["SFR"]["sub-mm"],
        # type_dict['SFR']['IR'],
        # type_dict['SFR']['UV'] + type_dict['SFR']['Vis'],
    ]

    sfr_colors = [
        "dodgerblue",
        "darkorange",
        "forestgreen",
        #'black',
        #'violet',
    ]

    types = [
        "cm",
        "mm",
        "sub-mm",
        "IR",
        "UV/Vis",
    ]

    def make_labels(list):
        new_labels = []
        total = sum(list)
        for i in list:
            percent = 100 * i / total
            new_labels.append(r"{:.1f}\%".format(percent))
        return new_labels

    axs[0, 0].pie(
        carbon_data,
        labels=make_labels(carbon_data),
        colors=carbon_colors,
        labeldistance=1.1,
        wedgeprops={"linewidth": 1.0, "edgecolor": "black", "alpha": 0.5},
    )
    axs[0, 0].annotate(
        r"\textbf{Carbon Stars}",
        xy=(0.5, 1.05),
        xycoords="axes fraction",
        color="black",
        ha="center",
        va="top",
        size=30,
    )

    axs[0, 1].pie(
        dark_data,
        labels=make_labels(dark_data),
        colors=dark_colors,
        labeldistance=1.1,
        wedgeprops={"linewidth": 1.0, "edgecolor": "black", "alpha": 0.5},
    )
    axs[0, 1].annotate(
        r"\textbf{Dark Clouds}", xy=(0.5, 1.05), xycoords="axes fraction", color="black", ha="center", va="top", size=30
    )

    wedges, _ = axs[1, 0].pie(
        los_data,
        labels=make_labels(los_data),
        colors=los_colors,
        labeldistance=1.1,
        wedgeprops={"linewidth": 1.0, "edgecolor": "black", "alpha": 0.5},
    )
    axs[1, 0].annotate(
        r"\textbf{LOS Clouds}", xy=(0.5, 1.05), xycoords="axes fraction", color="black", ha="center", va="top", size=30
    )

    axs[1, 1].pie(
        sfr_data,
        labels=make_labels(sfr_data),
        colors=sfr_colors,
        labeldistance=1.1,
        wedgeprops={"linewidth": 1.0, "edgecolor": "black", "alpha": 0.5},
    )
    axs[1, 1].annotate(
        r"\textbf{SFRs}", xy=(0.5, 1.05), xycoords="axes fraction", color="black", ha="center", va="top", size=30
    )

    axs[0, 1].legend(wedges, types, title="Wavelengths", loc="center left", bbox_to_anchor=(1, 0, 0.5, 1))

    fig.tight_layout()
    fig.subplots_adjust(wspace=0.1, hspace=0.1)
    plt.show()

    plt.savefig(
        filename if filename is not None else "waves_by_source_type.pdf",
        format="pdf",
        transparent=True,
        bbox_inches="tight",
    )

    return


def kappas(mol_list=None, nbins=100, filename=None):
    """
    Makes a histogram plot of kappa (Ray's Asymmetry Parameter) values.

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)   
    nbins : int
        The number of bins to use in the histogram (default is 100)  
    filename : str
        The filename for the output images (default is 'kappas.pdf')
    """

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    # gather the data
    kappas = [x.kappa for x in mol_list if x.kappa]

    # get the plot set up
    plt.close("Kappas")
    fig = plt.figure(num="Kappas", figsize=(20, 8))
    plt.ion()

    # set some font defaults
    fontparams = {"size": 24, "family": "sans-serif", "sans-serif": ["Helvetica"]}
    plt.rc("font", **fontparams)
    plt.rc("mathtext", fontset="stixsans")

    # load up an axis
    ax = fig.add_subplot(111)
    ax.tick_params(axis="x", which="both", direction="in", length=5, width=1)
    ax.tick_params(axis="y", which="both", direction="in", length=5, width=1)

    # label the axes
    plt.xlabel(r"$\kappa$")
    plt.ylabel(r"\# Molecules")

    ax.hist(kappas, nbins, facecolor="dodgerblue", alpha=0.25)
    ax.hist(kappas, nbins, edgecolor="dodgerblue", linewidth=1.5, fill=False)

    # fix the ticks
    ax.set_xticks([-1.0, -0.5, 0.0, 0.5, 1.0])
    ax.set_yscale("log")
    ax.yaxis.set_major_formatter(ScalarFormatter())

    plt.show()

    plt.savefig(
        filename if filename is not None else "kappas.pdf",
        format="pdf",
        transparent=True,
        bbox_inches="tight",
        pad_inches=0,
    )

    return


def waves_pie_chart(mol_list=None, filename=None):
    """
    Makes a pie chart of the fraction of interstellar molecules that are 
    detected a cm, mm, sub-mm, IR, UV, and Vis wavelengths.
{"cm": 0, "mm": 0, "sub-mm": 0, "IR": 0, "UV": 0, "Vis": 0},
    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)
    filename : str
        The filename for the output images (default is 'type_pie_chart.pdf')
    """

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    # Close an old figure if it exists and initialize a new figure
    plt.close("Waves Pie Chart")
    plt.figure(num="Waves Pie Chart", figsize=(10, 8))
    plt.ion()

    # set some font defaults
    fontparams = {"size": 24, "family": "sans-serif", "sans-serif": ["Helvetica"]}
    plt.rc("font", **fontparams)
    plt.rc("mathtext", fontset="stixsans")

    # gather the data
    my_dict = {
		"radio" : 0,
        "IR": 0,
        "UV/Vis": 0,
    }

    for mol in mol_list:
        if "cm" in mol.wavelengths or "mm" in mol.wavelengths or "sub-mm" in mol.wavelengths:
            my_dict["radio"] += 1
        if "IR" in mol.wavelengths:
            my_dict["IR"] += 1                   
        if "UV" in mol.wavelengths or "Vis" in mol.wavelengths:
            my_dict["UV/Vis"] += 1
            
    nmols = len(mol_list)
    for type in my_dict:
        my_dict[type] = my_dict[type] / nmols

    labels = ["radio", "IR", "UV/Vis"]
    vals = [my_dict[x] for x in labels]
    fracs = [my_dict[x] for x in labels]
    colors = ["dodgerblue", "violet", "red"]

    # set up a plot
    ax = plt.subplot(111)

    
    plt.pie(vals, labels = labels, colors=colors, autopct='%1.1f\%%', 
        pctdistance=0.8,
        wedgeprops = {"edgecolor" : "black",
                      'linewidth': 1, 'alpha':0.7})
    
    
#     size = 0.1
# 
#     def getshift(x):
#         return -(90 - (360 - 360 * x) / 2)
# 
#     ax.pie(
#         [fracs[0], 1.0 - fracs[0]],
#         colors=["dodgerblue", "#EEEEEE"],
#         radius=1,
#         startangle=getshift(fracs[0]),
#         wedgeprops=dict(width=size, edgecolor="w", linewidth=1),
#     )
#     ax.pie(
#         [fracs[1], 1.0 - fracs[1]],
#         colors=["darkorange", "#EEEEEE"],
#         radius=1 - size - 0.02,
#         startangle=getshift(fracs[1]),
#         wedgeprops=dict(width=size, edgecolor="w", linewidth=1),
#     )
#     ax.pie(
#         [fracs[2], 1.0 - fracs[2]],
#         colors=["violet", "#EEEEEE"],
#         radius=1 - 2 * size - 0.06,
#         startangle=getshift(fracs[2]),
#         wedgeprops=dict(width=size, edgecolor="w", linewidth=1),
#     )
# 
#     ax.annotate(r"\textbf{radio}", xy=(0.5, 0.11), xycoords="axes fraction", color="dodgerblue", ha="center", size=14)
#     ax.annotate(r"\textbf{IR}", xy=(0.5, 0.16), xycoords="axes fraction", color="darkorange", ha="center", size=14)
#     ax.annotate(r"\textbf{UV/Vis}", xy=(0.5, 0.255), xycoords="axes fraction", color="violet", ha="center", size=14)
# 
#     percents = [r"\textbf{" + "{:.1f}".format((x * 100)) + r"}\%" for x in fracs]
# 
#     start = 0.585
#     shift = 0.0485
#     ax.annotate(
#         percents[0],
#         xy=(start + 6 * shift, 0.5),
#         xycoords="axes fraction",
#         color="white",
#         ha="center",
#         va="center",
#         size=12,
#         rotation=-90,
#     )
#     ax.annotate(
#         percents[1],
#         xy=(start + 5 * shift, 0.5),
#         xycoords="axes fraction",
#         color="darkorange",
#         ha="center",
#         va="center",
#         size=12,
#         rotation=-90,
#     )
#     ax.annotate(
#         percents[2],
#         xy=(start + 4 * shift, 0.5),
#         xycoords="axes fraction",
#         color="forestgreen",
#         ha="center",
#         va="center",
#         size=12,
#         rotation=-90,
#     )

    plt.tight_layout()


    plt.savefig(
        filename if filename is not None else "waves_pie_chart.png", format="png", transparent=True, bbox_inches="tight", dpi=600
    )  # ,pad_inches=-.65)
    
    plt.show()

    return 

#############################################################
# 						    LaTeX	 						#
#############################################################


def make_ism_tables(mol_list=None, filename=None):
    """
    Generates the two latex tables for the census that contain all the 
    detected ISM/CSM molecules, using the hyperlink tags for the census paper.
    
    The tables are broken out by number of atoms, with the first table containing
    two-atom molecules through seven-atom molecules and the second table containing
    molecules with eight or more atoms.  It outputs two files, following the convetion 
    of filename + '_2-7.tex' and filename + '_8+.tex'

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)   
    filename : str
        The filename for the output images (default is 'ism_table')
    """

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    # If output filename is not specified, use defaults
    if filename is None:
        filename_two_seven = "ism_table_2-7.tex"
        filename_eight_more = "ism_table_8+.tex"
    else:
        filename_two_seven = filename + "_2-7.tex"
        filename_eight_more = filename + "_8+.tex"

    # put molecules into individual lists for ease of use below
    two_atoms = [x for x in mol_list if x.natoms == 2]
    three_atoms = [x for x in mol_list if x.natoms == 3]
    four_atoms = [x for x in mol_list if x.natoms == 4]
    five_atoms = [x for x in mol_list if x.natoms == 5]
    six_atoms = [x for x in mol_list if x.natoms == 6]
    seven_atoms = [x for x in mol_list if x.natoms == 7]
    eight_atoms = [x for x in mol_list if x.natoms == 8]
    nine_atoms = [x for x in mol_list if x.natoms == 9]
    ten_atoms = [x for x in mol_list if x.natoms == 10]
    eleven_atoms = [x for x in mol_list if x.natoms == 11]
    twelve_atoms = [x for x in mol_list if x.natoms == 12]
    thirteen_atoms = [x for x in mol_list if x.natoms == 13]
    pahs = [x for x in mol_list if x.pah is True]
    fullerenes = [x for x in mol_list if x.fullerene is True]

    # we'll do the first table with molecules from 2-7 atoms in it.

    # initialize a list to hold each line of the table
    table_list = []

    # Some of our columns will be double: two, three, four, and five atoms. So we need to split those into sub-lists
    i = int(len(two_atoms) / 2) + 1
    j = int(len(three_atoms) / 2) + 1
    k = int(len(four_atoms) / 2) + 1
    n = int(len(five_atoms) / 2) + 1

    two_seven = [
        two_atoms[:i],
        two_atoms[i:],
        three_atoms[:j],
        three_atoms[j:],
        four_atoms[:k],
        four_atoms[k:],
        five_atoms[:n],
        five_atoms[n:],
        six_atoms,
        seven_atoms,
    ]

    # figure out the max number of rows for our table
    nlines = np.max([len(x) for x in two_seven])

    # put into latex and add to the table_list

    for x in range(nlines):
        table_line = ""
        for r in two_seven:
            if x < len(r):
                label = r[x].label
                formula = r[x].formula if r[x].table_formula is None else r[x].table_formula
                table_line += "\\hyperref[{}]{{\ce{{{}}}}}\t&\t".format(label, formula)
            else:
                table_line += "\t&\t"
        table_line = table_line[:-2]
        table_line += "\\\\\n"
        table_list.append(table_line)

    with open(filename_two_seven, "w") as output:

        output.write("\\begin{table*}\n")
        output.write("\\centering\n")
        output.write(
            "\\caption{List of detected interstellar molecules with two to seven atoms, categorized by number of atoms, and vertically ordered by detection year.  Column headers and molecule formulas are in-document hyperlinks in most PDF viewers.}\n"
        )
        output.write(
            "\\begin{tabular*}{\\textwidth}{l l @{\\extracolsep{\\fill}} l l  @{\\extracolsep{\\fill}} l l  @{\\extracolsep{\\fill}} l l @{\\extracolsep{\\fill}} l @{\\extracolsep{\\fill}} l}\n"
        )
        output.write("\\hline\\hline\n")
        output.write(
            "\\multicolumn{2}{c}{\\hyperref[2atoms]{2 Atoms}} &\multicolumn{2}{c}{\\hyperref[3atoms]{3 Atoms}}& \multicolumn{2}{c}{\\hyperref[4atoms]{4 Atoms}} & \multicolumn{2}{c}{\\hyperref[5atoms]{5 Atoms}} & \\hyperref[6atoms]{6 Atoms} & \\hyperref[7atoms]{7 Atoms} \\\\\n"
        )
        output.write("\\hline\n")

        for x in table_list:
            output.write(x)

        output.write("\\hline\n\\end{tabular*}\n\\label{two_seven}\n\\end{table*}\endinput")

    # now we do the 8 or more table.
    table_list = []

    eight_more = [
        eight_atoms,
        nine_atoms,
        ten_atoms,
        eleven_atoms,
        twelve_atoms,
        thirteen_atoms,
        pahs,
        fullerenes,
    ]

    nlines = np.max([len(x) for x in eight_more])

    for x in range(nlines):
        table_line = ""
        for r in eight_more:
            if x < len(r):
                label = r[x].label
                formula = r[x].formula if r[x].table_formula is None else r[x].table_formula
                table_line += "\\hyperref[{}]{{\ce{{{}}}}}\t&\t".format(label, formula)
            else:
                table_line += "\t&\t"
        table_line = table_line[:-2]
        table_line += "\\\\\n"
        table_list.append(table_line)

    with open(filename_eight_more, "w") as output:

        output.write("\\begin{table*}\n")
        output.write("\\centering\n")
        output.write(
            "\\caption{List of detected interstellar molecules with eight or more atoms, categorized by number of atoms, and vertically ordered by detection year.  Column headers and molecule formulas are in-document hyperlinks in most PDF viewers.}\n"
        )
        output.write(
            "\\begin{tabular*}{\\textwidth}{l @{\\extracolsep{\\fill}} l @{\\extracolsep{\\fill}} l @{\\extracolsep{\\fill}} l @{\\extracolsep{\\fill}} l @{\\extracolsep{\\fill}} l @{\\extracolsep{\\fill}} l @{\\extracolsep{\\fill}} l }\n"
        )
        output.write("\\hline\\hline\n")
        output.write(
            "\\hyperref[8atoms]{8 Atoms} & \\hyperref[9atoms]{9 Atoms} & \\hyperref[10atoms]{10 Atoms} & \\hyperref[11atoms]{11 Atoms} & \\hyperref[12atoms]{12 Atoms} & \\hyperref[13atoms]{13 Atoms} & \\hyperref[pahs]{PAHs} & \\hyperref[fullerenes]{Fullerenes}  \\\\\n"
        )
        output.write("\\hline\n")

        for x in table_list:
            output.write(x)

        output.write("\\hline\n\\end{tabular*}\n\\label{eight_more}\n\\end{table*}\endinput")

def make_exgal_count(mol_list=None, filename=None):
    """
    Makes a .tex file containing the number of molecules in mol_list found in external galaxies

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)   
    filename : str
        The filename for the output snippet (default is 'nexgal.tex')
    """

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    with open(filename if filename is not None else "nexgal.tex", "w") as output:
        output.write(f"{len([x for x in mol_list if x.exgal is True])}\endinput")

def make_exgal_percent(mol_list=None, filename=None):
    """
    Makes a .tex file containing the percentage of molecules in mol_list found in external galaxies

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)   
    filename : str
        The filename for the output snippet (default is 'nexgalpercent.tex')
    """

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    with open(filename if filename is not None else "nexgalpercent.tex", "w") as output:
        output.write(f"{100*len([x for x in mol_list if x.exgal is True])/len(mol_list):.0f}\endinput")

def _make_exgal_sentence():
    """
    Makes a .tex file containing a sentence for the census with top exgal detection sources.

    Not recommended for any other uses.
    """

    ngc253 = 0
    pks = 0
    m82 = 0

    for x in all_molecules:
        if x.exgal is True:
            if 'NGC 253' in x.exgal_sources:
                ngc253 +=1
            if 'PKS 1830' in x.exgal_sources:
                pks +=1
            if 'M82' in x.exgal_sources:
                m82 +=1

    my_str = f"the line of sight to PKS 1830-211 ({pks} molecules), NGC 253 ({ngc253} molecules), and M82 ({m82} molecules)\endinput"

    with open("exgal_sentence.tex", "w") as output:
        output.write(my_str)

def make_det_count(mol_list=None, filename=None):
    """
    Makes a .tex file containing the number of molecules in mol_list

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)   
    filename : str
        The filename for the output snippet (default is 'ndetects.tex')
    """

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    with open(filename if filename is not None else "ndetects.tex", "w") as output:
        output.write(f"{len(mol_list)}\endinput")


def make_elem_count(mol_list=None, filename=None):
    """
    Makes a .tex file containing the number of unique elements in molecules in mol_list

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)   
    filename : str
        The filename for the output snippet (default is 'nelems.tex')
    """

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    elems = []

    for mol in mol_list:
        for atom in mol.atoms:
            if mol.atoms[atom] > 0:
                if atom not in elems:
                    # we don't want to count deuterium in the NH3D+ as special
                    if atom != 'D':
                        elems.append(atom)

    with open(filename if filename is not None else "nelems.tex", "w") as output:
        output.write(f"{len(elems)}\endinput")

def make_ppd_count(mol_list=None, filename=None):
    """
    Makes a .tex file containing the number of molecules from mol_list 
    found in protoplanetary disks. Does not count isotopologues, which are 
    tabulated separately.

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)   
    filename : str
        The filename for the output snippet (default is 'nppds.tex')
    """ 

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    ppd_mols = 0

    for mol in mol_list:
        if mol.ppd is True:
            ppd_mols += 1

    with open(filename if filename is not None else "nppds.tex", "w") as output:
        output.write(f"{ppd_mols}\endinput")        

def make_exo_count(mol_list=None, filename=None):
    """
    Makes a .tex file containing the number of molecules from mol_list 
    found in exoplanetary atmospheres. 

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)   
    filename : str
        The filename for the output snippet (default is 'nexos.tex')
    """ 

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    exo_mols = 0

    for mol in mol_list:
        if mol.exo is True:
            exo_mols += 1

    with open(filename if filename is not None else "nexos.tex", "w") as output:
        output.write(f"{exo_mols}\endinput")                

def make_ices_count(mol_list=None, filename=None):
    """
    Makes a .tex file containing the number of molecules from mol_list 
    found in ices. 

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)   
    filename : str
        The filename for the output snippet (default is 'nices.tex')
    """ 

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    ice_mols = 0

    for mol in mol_list:
        if mol.ice is True:
            ice_mols += 1

    with open(filename if filename is not None else "nices.tex", "w") as output:
        output.write(f"{ice_mols}\endinput")                

def make_ppd_isos_count(mol_list=None, filename=None):
    """
    Makes a .tex file containing the number of isotopologues of molecules from 
    mol_list found in protoplanetary disks.

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)   
    filename : str
        The filename for the output snippet (default is 'nppdisos.tex')
    """ 

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    ppd_iso_mols = 0

    for mol in mol_list:
        if mol.ppd is True:
            if mol.ppd_isos is not None:
                ppd_iso_mols += len(mol.ppd_isos)

    with open(filename if filename is not None else "nppdisos.tex", "w") as output:
        output.write(f"{ppd_iso_mols}\endinput")               


def make_exgal_table(mol_list=None, filename=None):
    """
    Generates the latex table for the census that contains all the detected 
    extragalactic molecules, using the provided bibtex references.

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)   
    filename : str
        The filename for the output file (default is 'exgal_table.tex')
    """ 

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    # put molecules into individual lists for ease of use below
    two_atoms = [x for x in mol_list if x.natoms == 2 and x.exgal is not None]
    three_atoms = [x for x in mol_list if x.natoms == 3 and x.exgal is not None]
    four_atoms = [x for x in mol_list if x.natoms == 4 and x.exgal is not None]
    five_atoms = [x for x in mol_list if x.natoms == 5 and x.exgal is not None]
    six_atoms = [x for x in mol_list if x.natoms == 6 and x.exgal is not None]
    seven_atoms = [x for x in mol_list if x.natoms == 7 and x.exgal is not None]
    eight_atoms = [x for x in mol_list if x.natoms == 8 and x.exgal is not None]
    nine_atoms = [x for x in mol_list if x.natoms == 9 and x.exgal is not None]
    # ten_atoms = [x for x in mol_list if x.natoms == 10 and x.exgal is not None] # no hits here yet
    # eleven_atoms = [x for x in mol_list if x.natoms == 11 and x.exgal is not None] # no hits here yet
    twelve_atoms = [x for x in mol_list if x.natoms == 12 and x.exgal is not None]

    # initialize a list to hold each line of the first table
    table = []

    # add the pre-amble latex stuff
    table.append(r"\begin{table*}" + "\n")
    table.append(r"\centering" + "\n")
    table.append(
        r"\caption{List of molecules detected in external galaxies with references to the first detections.  Tentative detections are indicated, and some extra references are occasionally provided for context.}"
        + "\n"
    )
    table.append(
        r"\begin{tabular*}{\textwidth}{l @{\extracolsep{\fill}} l @{\extracolsep{\fill}} l @{\extracolsep{\fill}} l @{\extracolsep{\fill}} l @{\extracolsep{\fill}} l @{\extracolsep{\fill}} l @{\extracolsep{\fill}} l @{\extracolsep{\fill}} l @{\extracolsep{\fill}} l @{\extracolsep{\fill}}}"
        + "\n"
    )
    table.append(r"\hline\hline" + "\n")
    table.append(
        r"\multicolumn{2}{c}{2 Atoms}&\multicolumn{2}{c}{3 Atoms}&\multicolumn{2}{c}{4 Atoms}&\multicolumn{2}{c}{5 Atoms}\\"
        + "\n"
    )
    table.append(r"Species	&	Ref.	&	Species	&	Ref.	&	Species	&	Ref.	&	Species	&	Ref.	\\" + "\n")
    table.append(r"\hline" + "\n")

    # initialize a dictionary to hold references and the Ref ID that is used, and initialize a counter
    refs_dict = {}
    ref_idx = 1

    two_five = [
        two_atoms,
        three_atoms,
        four_atoms,
        five_atoms,
    ]

    six_more = [
        six_atoms,
        seven_atoms,
        eight_atoms,
        nine_atoms,
        twelve_atoms,
    ]

    # Do the two to five atom table first

    # figure out the max number of rows for the first half of our table
    nlines_first = np.max([len(x) for x in two_five])

    # now just loop through and add
    for i in range(nlines_first):
        table_line = ""
        for x in two_five:
            if i < len(x):
                # get the formula in there
                formula = x[i].formula if x[i].table_formula is None else x[i].table_formula
                # if it's not tentative, just add it
                if x[i].exgal is True:
                    table_line += f"\ce{{{formula}}}\t&\t"
                # if it is tentative, note that
                if x[i].exgal == "Tentative":
                    table_line += f"\ce{{{formula}}}" + r"$^{\dagger}$" + "\t&\t"
                # now we deal with the references
                ref_strs = []
                # some detections have more than one reference, so we loop over them
                for ref_ID in x[i].exgal_d_bib_ids:
                    # if the reference has already been used before, we just use the number from that instance
                    if ref_ID in refs_dict:
                        ref_strs.append(refs_dict[ref_ID])
                    # if it hasn't been used before
                    else:
                        # assign it the next number in line to be used
                        ref_strs.append(str(ref_idx))
                        # add it as an entry in refs_dict so we can use it again if needs be
                        refs_dict[ref_ID] = str(ref_idx)
                        # and increment the number
                        ref_idx += 1
                # join the reference string together and add it to the line
                table_line += f"{', '.join(ref_strs)}\t&\t"
            else:
                table_line += "\t&\t&\t"
        table_line = table_line[:-2]
        table_line += "\\\\\n"
        table.append(table_line)

    # add stuff between the two halves of the table
    table.append(r"\hline\hline" + "\n")
    table.append(r"\end{tabular*}" + "\n")
    table.append(
        r"\begin{tabular*}{\textwidth}{l @{\extracolsep{\fill}} l @{\extracolsep{\fill}} l @{\extracolsep{\fill}} l @{\extracolsep{\fill}} l @{\extracolsep{\fill}} l @{\extracolsep{\fill}} l @{\extracolsep{\fill}} l @{\extracolsep{\fill}} l @{\extracolsep{\fill}} l @{\extracolsep{\fill}} l @{\extracolsep{\fill}} l @{\extracolsep{\fill}}}"
        + "\n"
    )
    table.append(
        r"\multicolumn{2}{c}{6 Atoms}&\multicolumn{2}{c}{7 Atoms}&\multicolumn{2}{c}{8 Atoms}&\multicolumn{2}{c}{9 Atoms}&\multicolumn{2}{c}{12 Atoms}\\"
        + "\n"
    )
    table.append(r"Species	&	Ref.	&	Species	&	Ref.	&	Species	&	Ref.	&	Species	&	Ref.	&	Species	&	Ref.	\\" + "\n")
    table.append(r"\hline" + "\n")

    # figure out the max number of rows for the second half of our table
    nlines_second = np.max([len(x) for x in six_more])

    # now just loop through and add
    for i in range(nlines_second):
        table_line = ""
        for x in six_more:
            if i < len(x):
                # get the formula in there
                formula = x[i].formula if x[i].table_formula is None else x[i].table_formula
                # if it's not tentative, just add it
                if x[i].exgal is True:
                    table_line += f"\ce{{{formula}}}\t&\t"
                # if it is tentative, note that
                if x[i].exgal == "Tentative":
                    table_line += f"\ce{{{formula}}}" + r"$^{\dagger}$" + "\t&\t"
                # now we deal with the references
                ref_strs = []
                # some detections have more than one reference, so we loop over them
                for ref_ID in x[i].exgal_d_bib_ids:
                    # if the reference has already been used before, we just use the number from that instance
                    if ref_ID in refs_dict:
                        ref_strs.append(refs_dict[ref_ID])
                    # if it hasn't been used before
                    else:
                        # assign it the next number in line to be used
                        ref_strs.append(str(ref_idx))
                        # add it as an entry in refs_dict so we can use it again if needs be
                        refs_dict[ref_ID] = str(ref_idx)
                        # and increment the number
                        ref_idx += 1
                # join the reference string together and add it to the line
                table_line += f"{', '.join(ref_strs)}\t&\t"
            else:
                table_line += "\t&\t&\t"
        table_line = table_line[:-2]
        table_line += "\\\\\n"
        table.append(table_line)

    # close out the data portion of the table
    table.append(r"\hline" + "\n")
    table.append(r"\end{tabular*}" + "\n")
    table.append(r"\justify" + "\n")
    table.append(r"$^{\dagger}$Tentative detection\\" + "\n")

    # now we do the references; we'll have to play a few tricks to get an ordered list out of the dictionary
    reference_ids = []  # the actual latex citekeys
    reference_idx = []  # the index number from the table
    for ref in refs_dict:
        reference_ids.append(ref)
        reference_idx.append(int(refs_dict[ref]))
    # get the ordered array that will sort the reference indices
    sort_idx = np.argsort(np.array(reference_idx))
    # then use it to sort both the reference indices and the references themselves
    reference_ids = np.array(reference_ids)[sort_idx]
    reference_idx = np.array(reference_idx)[sort_idx]

    # get the preamble text out of the way; be sure not to include a carriage return
    table.append(r"\textbf{References:}")

    # figure out how many references we have, and iterate over them
    for i in range(np.max(reference_idx)):
        table.append(f" [{reference_idx[i]}] " + r"\citet{" + f"{reference_ids[i]}" + r"} ")

    # add the closing bit and a carriage return and the last lines
    table.append(r"\\" + "\n")
    table.append(r"\label{exgal_mols}" + "\n")
    table.append(r"\end{table*}\endinput")

    with open(filename if filename is not None else "exgal_table.tex", "w") as output:
        for line in table:
            output.write(line)

    return


def make_ppd_table(mol_list=None, filename=None):
    """
    Generates the latex table for the census that contains all the detected 
    protoplanetary disk molecules, using the provided bibtex references.

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)   
    filename : str
        The filename for the output file (default is 'ppd_table.tex')
    """ 

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    # put molecules into individual lists for ease of use below; pull out isotopologues if they exist
    two_atoms = []
    three_atoms = []
    four_atoms = []
    five_atoms = []
    six_atoms = []

    for mol in mol_list:
        if mol.natoms == 2:
            if mol.ppd is True:
                two_atoms.append(mol)
            if mol.ppd_isos is not None:
                for iso in mol.ppd_isos:
                    two_atoms.append(iso)
        if mol.natoms == 3:
            if mol.ppd is True:
                three_atoms.append(mol)
            if mol.ppd_isos is not None:
                for iso in mol.ppd_isos:
                    three_atoms.append(iso)
        if mol.natoms == 4:
            if mol.ppd is True:
                four_atoms.append(mol)
            if mol.ppd_isos is not None:
                for iso in mol.ppd_isos:
                    four_atoms.append(iso)
        if mol.natoms == 5:
            if mol.ppd is True:
                five_atoms.append(mol)
            if mol.ppd_isos is not None:
                for iso in mol.ppd_isos:
                    five_atoms.append(iso)
        if mol.natoms == 6:
            if mol.ppd is True:
                six_atoms.append(mol)
            if mol.ppd_isos is not None:
                for iso in mol.ppd_isos:
                    six_atoms.append(iso)

    # initialize a list to hold each line of the first table
    table = []

    # add the pre-amble latex stuff
    table.append(r"\begin{table*}" + "\n")
    table.append(r"\centering" + "\n")
    table.append(
        r"\caption{List of molecules, including rare isotopic species, detected in protoplanetary disks, with references to representative detections.  The earliest reported detection of a species in the literature is provided on a best-effort basis.  Tentative and disputed detections are not included (see text).}"
        + "\n"
    )
    table.append(
        r"\begin{tabular*}{\textwidth}{l @{\extracolsep{\fill}} l @{\extracolsep{\fill}} l @{\extracolsep{\fill}} l @{\extracolsep{\fill}} l @{\extracolsep{\fill}} l @{\extracolsep{\fill}} l @{\extracolsep{\fill}} l @{\extracolsep{\fill}} l @{\extracolsep{\fill}} l @{\extracolsep{\fill}}}"
        + "\n"
    )
    table.append(r"\hline\hline" + "\n")
    table.append(
        r"\multicolumn{2}{c}{2 Atoms}&\multicolumn{2}{c}{3 Atoms}&\multicolumn{2}{c}{4 Atoms}&\multicolumn{2}{c}{5 Atoms}&\multicolumn{2}{c}{6 Atoms}\\"
        + "\n"
    )
    table.append(r"Species	&	Ref.	&	Species	&	Ref.	&	Species	&	Ref.	&	Species	&	Ref.	&	Species	&	Ref.	\\" + "\n")
    table.append(r"\hline" + "\n")

    # initialize a dictionary to hold references and the Ref ID that is used, and initialize a counter
    refs_dict = {}
    ref_idx = 1

    two_six = [
        two_atoms,
        three_atoms,
        four_atoms,
        five_atoms,
        six_atoms,
    ]

    # figure out the max number of rows for the table
    nlines = np.max([len(x) for x in two_six])

    # now just loop through and add
    for i in range(nlines):
        table_line = ""
        for x in two_six:
            if i < len(x):
                # get the formula in there
                formula = x[i].formula if x[i].table_formula is None else x[i].table_formula
                # if it's not tentative, just add it
                table_line += f"\ce{{{formula}}}\t&\t"
                # now we deal with the references
                ref_strs = []
                # some detections have more than one reference, so we loop over them
                for ref_ID in x[i].ppd_d_bib_ids:
                    # if the reference has already been used before, we just use the number from that instance
                    if ref_ID in refs_dict:
                        ref_strs.append(refs_dict[ref_ID])
                    # if it hasn't been used before
                    else:
                        # assign it the next number in line to be used
                        ref_strs.append(str(ref_idx))
                        # add it as an entry in refs_dict so we can use it again if needs be
                        refs_dict[ref_ID] = str(ref_idx)
                        # and increment the number
                        ref_idx += 1
                # join the reference string together and add it to the line
                table_line += f"{', '.join(ref_strs)}\t&\t"
            else:
                table_line += "\t&\t&\t"
        table_line = table_line[:-2]
        table_line += "\\\\\n"
        table.append(table_line)

    # close out the data portion of the table
    table.append(r"\hline" + "\n")
    table.append(r"\end{tabular*}" + "\n")
    table.append(r"\justify" + "\n")

    # now we do the references; we'll have to play a few tricks to get an ordered list out of the dictionary
    reference_ids = []  # the actual latex citekeys
    reference_idx = []  # the index number from the table
    for ref in refs_dict:
        reference_ids.append(ref)
        reference_idx.append(int(refs_dict[ref]))
    # get the ordered array that will sort the reference indices
    sort_idx = np.argsort(np.array(reference_idx))
    # then use it to sort both the reference indices and the references themselves
    reference_ids = np.array(reference_ids)[sort_idx]
    reference_idx = np.array(reference_idx)[sort_idx]

    # get the preamble text out of the way; be sure not to include a carriage return
    table.append(r"\textbf{References:}")

    # figure out how many references we have, and iterate over them
    for i in range(np.max(reference_idx)):
        table.append(f" [{reference_idx[i]}] " + r"\citet{" + f"{reference_ids[i]}" + r"} ")

    # add the closing bit and a carriage return and the last lines
    table.append(r"\\" + "\n")
    table.append(r"\label{ppd_mols}" + "\n")
    table.append(r"\end{table*}\endinput")

    with open(filename if filename is not None else "ppd_table.tex", "w") as output:
        for line in table:
            output.write(line)

    return


def make_exo_table(mol_list=None, filename=None):
    """
    Generates the latex table for the census that contains all the detected 
    exoplanetary atmosphere molecules, using the provided bibtex references.

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)   
    filename : str
        The filename for the output file (default is 'exo_table.tex')
    """

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    # a list to hold detections
    detects = []

    for mol in mol_list:
        if mol.exo is True:
            detects.append(mol)

    # initialize a list to hold each line of the first table
    table = []

    # add the pre-amble latex stuff
    table.append(r"\begin{table}" + "\n")
    table.append(r"\centering" + "\n")
    table.append(
        r"\caption{List of molecules detected in exoplanetary atmospheres, with references to representative detections.  Tentative and disputed detections are not included.}"
        + "\n"
    )
    table.append(r"\begin{tabular*}{\columnwidth}{l @{\extracolsep{\fill}} l @{\extracolsep{\fill}}}" + "\n")
    table.append(r"\hline\hline" + "\n")
    table.append(r"Species & References\\" + "\n")
    table.append(r"\hline" + "\n")

    # initialize a dictionary to hold references and the Ref ID that is used, and initialize a counter
    refs_dict = {}
    ref_idx = 1

    # now just loop through and add
    for x in detects:
        table_line = ""
        # get the formula in there
        formula = x.formula if x.table_formula is None else x.table_formula
        table_line += f"\ce{{{formula}}}\t&\t"
        # now we deal with the references
        ref_strs = []
        # some detections have more than one reference, so we loop over them
        for ref_ID in x.exo_d_bib_ids:
            # if the reference has already been used before, we just use the number from that instance
            if ref_ID in refs_dict:
                ref_strs.append(refs_dict[ref_ID])
            # if it hasn't been used before
            else:
                # assign it the next number in line to be used
                ref_strs.append(str(ref_idx))
                # add it as an entry in refs_dict so we can use it again if needs be
                refs_dict[ref_ID] = str(ref_idx)
                # and increment the number
                ref_idx += 1
        # join the reference string together and add it to the line
        table_line += f"{', '.join(ref_strs)}" + r"\\" + "\n"
        # add it to the table.
        table.append(table_line)

    # close out the data portion of the table
    table.append(r"\hline" + "\n")
    table.append(r"\end{tabular*}" + "\n")
    table.append(r"\justify" + "\n")

    # now we do the references; we'll have to play a few tricks to get an ordered list out of the dictionary
    reference_ids = []  # the actual latex citekeys
    reference_idx = []  # the index number from the table
    for ref in refs_dict:
        reference_ids.append(ref)
        reference_idx.append(int(refs_dict[ref]))
    # get the ordered array that will sort the reference indices
    sort_idx = np.argsort(np.array(reference_idx))
    # then use it to sort both the reference indices and the references themselves
    reference_ids = np.array(reference_ids)[sort_idx]
    reference_idx = np.array(reference_idx)[sort_idx]

    # get the preamble text out of the way; be sure not to include a carriage return
    table.append(r"\textbf{References:}")

    # figure out how many references we have, and iterate over them
    for i in range(np.max(reference_idx)):
        table.append(f" [{reference_idx[i]}] " + r"\citet{" + f"{reference_ids[i]}" + r"} ")

    # add the closing bit and a carriage return and the last lines
    table.append(r"\\" + "\n")
    table.append(r"\label{exoplanet_mols}" + "\n")
    table.append(r"\end{table}\endinput")

    with open(filename if filename is not None else "exo_table.tex", "w") as output:
        for line in table:
            output.write(line)

    return

def make_ice_table(mol_list=None, filename=None):
    """
    Generates the latex table for the census that contains all the detected 
    interstellar ice molecules, using the provided bibtex references.

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)   
    filename : str
        The filename for the output file (default is 'ice_table.tex')
    """

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    # a list to hold detections
    detects = []

    for mol in mol_list:
        if mol.ice is True:
            detects.append(mol)

    # manually add OCN-, since it's not counted as an ISM molecule
    OCNm = Molecule(
        formula='OCN-',
        ice=True,
        ice_d_bib_ids=["2005A&A...441..249V"],
    )

    # gotta figure out where to insert it into the list
    last_3 = ''.join([str(x.natoms) for x in detects]).rindex('3')

    # then insert it
    detects.insert(last_3+1,OCNm)

    # initialize a list to hold each line of the first table
    table = []

    # add the pre-amble latex stuff
    table.append(r"\begin{table}" + "\n")
    table.append(r"\centering" + "\n")
    table.append(
        r"\caption{List of molecules detected in interstellar ices, with references to representative detections.}"
        + "\n"
    )
    table.append(r"\begin{tabular*}{\columnwidth}{l @{\extracolsep{\fill}} l @{\extracolsep{\fill}}}" + "\n")
    table.append(r"\hline\hline" + "\n")
    table.append(r"Species & References\\" + "\n")
    table.append(r"\hline" + "\n")

    # initialize a dictionary to hold references and the Ref ID that is used, and initialize a counter
    refs_dict = {}
    ref_idx = 1

    # now just loop through and add
    for x in detects:
        table_line = ""
        # get the formula in there
        formula = x.formula if x.table_formula is None else x.table_formula
        table_line += f"\ce{{{formula}}}\t&\t"
        # now we deal with the references
        ref_strs = []
        # some detections have more than one reference, so we loop over them
        for ref_ID in x.ice_d_bib_ids:
            # if the reference has already been used before, we just use the number from that instance
            if ref_ID in refs_dict:
                ref_strs.append(refs_dict[ref_ID])
            # if it hasn't been used before
            else:
                # assign it the next number in line to be used
                ref_strs.append(str(ref_idx))
                # add it as an entry in refs_dict so we can use it again if needs be
                refs_dict[ref_ID] = str(ref_idx)
                # and increment the number
                ref_idx += 1
        # join the reference string together and add it to the line
        table_line += f"{', '.join(ref_strs)}" + r"\\" + "\n"
        # add it to the table.
        table.append(table_line)

    # close out the data portion of the table
    table.append(r"\hline" + "\n")
    table.append(r"\end{tabular*}" + "\n")
    table.append(r"\justify" + "\n")

    # now we do the references; we'll have to play a few tricks to get an ordered list out of the dictionary
    reference_ids = []  # the actual latex citekeys
    reference_idx = []  # the index number from the table
    for ref in refs_dict:
        reference_ids.append(ref)
        reference_idx.append(int(refs_dict[ref]))
    # get the ordered array that will sort the reference indices
    sort_idx = np.argsort(np.array(reference_idx))
    # then use it to sort both the reference indices and the references themselves
    reference_ids = np.array(reference_ids)[sort_idx]
    reference_idx = np.array(reference_idx)[sort_idx]

    # get the preamble text out of the way; be sure not to include a carriage return
    table.append(r"\textbf{References:}")

    # figure out how many references we have, and iterate over them
    for i in range(np.max(reference_idx)):
        table.append(f" [{reference_idx[i]}] " + r"\citet{" + f"{reference_ids[i]}" + r"} ")

    # add the closing bit and a carriage return and the last lines
    table.append(r"\\" + "\n")
    table.append(r"\label{ice_mols}" + "\n")
    table.append(r"\end{table}\endinput")

    with open(filename if filename is not None else "ice_table.tex", "w") as output:
        for line in table:
            output.write(line)

    return

def make_det_per_year_by_atoms_table(mol_list=None, filename=None):
    """
    Generates the latex table for the census that contains the rates of 
    detectiosn of new molecules per year divided out by number of atoms.

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)   
    filename : str
        The filename for the output file (default is 'rates_by_atoms_table.tex')
    """

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    # a dictionary to hold the data.  Entries in the list ordered as follows:
    # 0 : start year for fitting (int)
    # 1 : array of years for plotting (np array)
    # 2 : number of cumulative detections for each year in that array (np array)
    # 3 : number of detections per year, if more than one year
    # 4 : R^2 value for the fit
    atoms_dict = {
        2: [1968],
        3: [1968],
        4: [1968],
        5: [1971],
        6: [1970],
        7: [1973],
        8: [1975],
        9: [1974],
        10: [2001],
        11: [2004],
        12: [2001],
        13: [2018],
        "PAHs": [2021],
        "Fullerenes": [2010],
    }

    # a function to count the cumulative detections to that point
    def _count_dets(mol_list, natoms, syear, eyear):
        dets = 0
        for mol in mol_list:
            if mol.natoms == natoms and syear <= mol.year <= eyear:
                dets += 1
        return dets

    # loop through the different categories and generate the years array and the associated detections.
    # the methods are inefficient, but so fast that it doesn't matter
    for x in atoms_dict:
        atoms_dict[x].append(np.arange(atoms_dict[x][0], date.today().year + 1))
        if isinstance(x, int):
            atoms_dict[x].append(np.array([_count_dets(mol_list, x, atoms_dict[x][0], z) for z in atoms_dict[x][1]]))
        elif x == "PAHs":
            atoms_dict[x].append([])
            for year in atoms_dict[x][1]:
                dets = 0
                for mol in mol_list:
                    if mol.pah is True and atoms_dict[x][0] <= mol.year <= year:
                        dets += 1
                atoms_dict[x][2].append(dets)
            atoms_dict[x][2] = np.array(atoms_dict[x][2])
        elif x == "Fullerenes":
            atoms_dict[x].append([])
            for year in atoms_dict[x][1]:
                dets = 0
                for mol in mol_list:
                    if mol.fullerene is True and atoms_dict[x][0] <= mol.year <= year:
                        dets += 1
                atoms_dict[x][2].append(dets)
            atoms_dict[x][2] = np.array(atoms_dict[x][2])

    # loop through and do a linear fit to get detections per year, if there's at least two years
    for x in atoms_dict:
        if len(atoms_dict[x][1]) > 1:
            slope, intercept, r_value, _, _ = linregress(atoms_dict[x][1], atoms_dict[x][2])
            atoms_dict[x].append(slope)
            atoms_dict[x].append(r_value)
        else:
            atoms_dict[x].append(None)
            atoms_dict[x].append(None)

    # initialize a list to hold each line of the first table
    table = []

    # add the pre-amble latex stuff
    table.append(r"\begin{table}[htb!]" + "\n")
    table.append(r"\centering" + "\n")
    table.append(
        r"\caption{Rates (\emph{m}) of detection of new molecules per year, sorted by number of atoms per molecule derived from linear fits to the data shown in Figure~\ref{cumulative_by_atoms} as well as the $R^2$ values of the fits, for molecules with 2--12 atoms.  The start year was chosen by the visual onset of a steady detection rate, and is given for each fit.  Rates and R$^2$ values obtained using the \texttt{scipy.stats.linregress} module.}"
        + "\n"
    )
    table.append(
        r"\begin{tabular*}{\columnwidth}{c @{\extracolsep{\fill}}  c @{\extracolsep{\fill}}  c @{\extracolsep{\fill}}  c }"
        + "\n"
    )
    table.append(r"\hline\hline" + "\n")
    table.append(r"\# Atoms    &   \emph{m} (yr$^{-1}$)    &   $R^2$       &   Onset Year      \\" + "\n")
    table.append(r"\hline" + "\n")

    # loop through and add lines to the table.
    for i in range(2, 14):
        if atoms_dict[i][3] is not None:
            table.append(
                f"{i}\t&\t{atoms_dict[i][3]:.2f}\t&\t{atoms_dict[i][4]:.2f}\t&\t{atoms_dict[i][0]}" + r"\\" + "\n"
            )
    if atoms_dict["PAHs"][3] is not None:
        table.append(
            f"PAHs\t&\t{atoms_dict['PAHs'][3]:.2f}\t&\t{atoms_dict['PAHs'][4]:.2f}\t&\t{atoms_dict['PAHs'][0]}"
            + r"\\"
            + "\n"
        )
    if atoms_dict["Fullerenes"] is not None:
        table.append(
            f"Fullerenes\t&\t{atoms_dict['Fullerenes'][3]:.2f}\t&\t{atoms_dict['Fullerenes'][4]:.2f}\t&\t{atoms_dict['Fullerenes'][0]}"
            + r"\\"
            + "\n"
        )

    # add the closing bit and a carriage return and the last lines
    table.append(r"\hline" + "\n")
    table.append(r"\end{tabular*}" + "\n")
    table.append(r"\label{rates_by_atoms_table}" + "\n")
    table.append(r"\end{table}\endinput")

    with open(filename if filename is not None else "rates_by_atoms_table.tex", "w") as output:
        for line in table:
            output.write(line)

    return

def make_facility_table(mol_list=None, filename=None):
    """
    Generates the latex table for the census that contains the number of 
    detections per facility.

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)   
    filename : str
        The filename for the output images (default is 'facilities_table.tex')
    """

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    # a dictionary to hold the data.
    facilities_dict = {}
    for x in all_telescopes:
        facilities_dict[x.shortname] = 0

    # a function to add up detections
    for mol in all_molecules:
        for scope in mol.telescopes:
            facilities_dict[scope.shortname] += 1

    # lists to then sort by
    scopes = []
    count = []

    for scope in facilities_dict:
        scopes.append(scope)
        count.append(facilities_dict[scope])

    scopes = np.array(scopes)
    count = np.array(count)
    sort_idx = np.argsort(count)[::-1]
    scopes = scopes[sort_idx]
    count = count[sort_idx]

    # calculate number of rows
    nrows = math.ceil(len(count) / 2)

    # initialize a list to hold each line of the first table
    table = []

    # add the pre-amble latex stuff
    table.append(r"\begin{table}[htb!]" + "\n")
    table.append(r"\centering" + "\n")
    table.append(r"\caption{Total number of detections for each facility listed in \S\ref{known}.}" + "\n")
    table.append(
        r"\begin{tabular*}{\columnwidth}{l @{\extracolsep{\fill}}  c @{\extracolsep{\fill}}  l @{\extracolsep{\fill}}  c }"
        + "\n"
    )
    table.append(r"\hline\hline" + "\n")
    table.append(r"Facility	&	\# 	&	Facility	&	\# \\" + "\n")
    table.append(r"\hline" + "\n")

    # loop through and add lines to the table.
    for i in range(nrows):
        if i + nrows < len(count):
            table.append(
                f"{scopes[i]}\t"
                + r"&"
                + f"\t{count[i]}\t"
                + r"&"
                + f"\t{scopes[i+nrows]}\t"
                + r"&"
                + f"\t{count[i+nrows]}\t"
                + r"\\"
                + "\n"
            )
        else:
            table.append(
                f"{scopes[i]}\t"
                + r"&"
                + f"\t{count[i]}\t"
                + r"&"
                + f"\t\t"
                + r"&"
                + f"\t\t"
                + r"\\"
                + "\n"
            )            

    # add the closing bit and a carriage return and the last lines
    table.append(r"\hline" + "\n")
    table.append(r"\end{tabular*}" + "\n")
    table.append(r"\label{detects_by_scope}" + "\n")
    table.append(r"\end{table}\endinput")

    with open(filename if filename is not None else "facilities_table.tex", "w") as output:
        for line in table:
            output.write(line)

    return


def make_source_table(mol_list=None, filename=None):
    """
    Generates the latex table for the census that contains the number of detections per source.

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)   
    filename : str
        The filename for the output file (default is 'source_table.tex')
    """

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    # a dictionary to hold the data.
    sources_dict = {}
    for x in all_sources:
        if x.type == "LOS Cloud":
            sources_dict["LOS Cloud"] = 0
        else:
            sources_dict[x.name] = 0

    # a function to add up detections
    for mol in all_molecules:
        for source in mol.sources:
            if source.type == "LOS Cloud":
                sources_dict["LOS Cloud"] += 1
            else:
                sources_dict[source.name] += 1

    # lists to then sort by
    sources = []
    count = []

    for src in sources_dict:
        sources.append(src)
        count.append(sources_dict[src])

    sources = np.array(sources)
    count = np.array(count)
    sort_idx = np.lexsort((sources,-count))
    sources = sources[sort_idx]
    count = count[sort_idx]

    # initialize a list to hold each line of the first table
    table = []

    # calculate number of rows
    nrows = math.ceil(len(count) / 2)    

    # add the pre-amble latex stuff
    table.append(r"\begin{table}[htb!]" + "\n")
    table.append(r"\centering" + "\n")
    table.append(
        r"\caption{Total number of detections that each source contributed to for the molecules listed in \S\ref{known}.  Detections made in clouds along the line of sight to a background source have been consolidated into `LOS Clouds,' and detections in closely-location regions have been group together as well (e.g. Sgr B2(OH), Sgr B2(N), Sgr B2(S), and Sgr B2(M) are all considered Sgr B2).}"
        + "\n"
    )
    table.append(
        r"\begin{tabular*}{\columnwidth}{l @{\extracolsep{\fill}}  c @{\extracolsep{\fill}}  l @{\extracolsep{\fill}}  c }"
        + "\n"
    )
    table.append(r"\hline\hline" + "\n")
    table.append(r"Source	&	\# 	&	Source	&	\# \\" + "\n")
    table.append(r"\hline" + "\n")

    # loop through and add lines to the table.
    for i in range(nrows):
        if i + nrows < len(count):
            table.append(
                f"{sources[i]}\t"
                + r"&"
                + f"\t{count[i]}\t"
                + r"&"
                + f"\t{sources[i+nrows]}\t"
                + r"&"
                + f"\t{count[i+nrows]}\t"
                + r"\\"
                + "\n"
            )
        else:
            table.append(
                f"{sources[i]}\t"
                + r"&"
                + f"\t{count[i]}\t"
                + r"&"
                + f"\t\t"
                + r"&"
                + f"\t\t"
                + r"\\"
                + "\n"
            )            

    # add the closing bit and a carriage return and the last lines
    table.append(r"\hline" + "\n")
    table.append(r"\end{tabular*}" + "\n")
    table.append(r"\label{detects_by_source}" + "\n")
    table.append(r"\end{table}\endinput")

    with open(filename if filename is not None else "source_table.tex", "w") as output:
        for line in table:
            output.write(line)

    return

def make_rate_counts(mol_list=None, syears=None, filename=None):
    """
    Makes .tex files containing the rate of new molecule detections of molecules in mol_list
    since the given years

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)
    syears : list
        A list of integers for the starting year to calculate trends (default is 1968 and 2005)   
    filename : str
        The based filename for the output snippets (default is 'rate_since_')
    """

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    # get the starting years, if they aren't set by the user
    if syears is None:
        # grab the earliest year in the list of molecules
        syears = [1968,2005]
    eyear = date.today().year

    # make the x-axis array of years
    years = np.arange(min([x.year for x in mol_list]), eyear + 1)

    # make an array to hold the detections
    dets = np.copy(years) * 0

    # loop through the years and the list and add everything up.  There's gotta be a better way to do this, but who cares, its fast enough.
    for x in range(len(dets)):
        i = 0
        for mol in mol_list:
            if mol.year < years[x] + 1:
                i += 1
        dets[x] = i

    # get some year indicies for years we care about
    def iyear(x):
        return np.argwhere(years == x)[0][0]

    # do linear fits to the data for the ranges we care about
    fit_results = {}
    for year in syears:
        fit_results[year] = np.polynomial.polynomial.Polynomial.fit(years[iyear(year) :], dets[iyear(year) :], 1).convert().coef[1]

    # write everything out
    for year in fit_results:
        with open(f'{filename}{year}.tex' if filename is not None else f'rate_since_{year}.tex', 'w') as output:
            output.write(f"{fit_results[year]:.1f}\endinput")

def make_percent_radio(mol_list=None, filename=None):
    """
    Makes a .tex file containing the percentage of molecules from mol_list 
    detected by radio astronomy. 

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)   
    filename : str
        The filename for the output snippet (default is 'radiopercent.tex')
    """ 

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    radio_mols = 0

    for mol in mol_list:
        if "cm" in mol.wavelengths or "mm" in mol.wavelengths or "sub-mm" in mol.wavelengths:
            radio_mols += 1

    with open(filename if filename is not None else "radiopercent.tex", "w") as output:
        output.write(f"{100*radio_mols/len(mol_list):.0f}\endinput") 

def make_scopes_count(telescopes_list=None, filename=None):
    """
    Makes a .tex file containing the number of telescopes used to detect ISM/CSM molecules

    Parameters
    ----------
    telescopes_list : list
        A list of telescope objects to use (default is all_telescopes)   
    filename : str
        The filename for the output snippet (default is 'nscopes.tex')
    """

    # If a list wasn't specified, default to all telescopes
    if telescopes_list is None:
        telescopes_list = all_telescopes

    i = 0
    # start the count
    for scope in telescopes_list:
        if scope.ndetects > 0:
            i += 1

    with open(filename if filename is not None else "nscopes.tex", "w") as output:
        output.write(f"{i}\endinput")

def make_percent_unsat(mol_list=None, filename=None):
    """
    Makes a .tex file containing the percentage of molecules from mol_list 
    that are hydrocarbons and at least partially unsaturated 

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)   
    filename : str
        The filename for the output snippet (default is 'unsatpercent.tex')
    """ 

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    un_sats = []
    for x in mol_list:
        if x.du is not None and x.du > 0.0:
            if 'H' in list(x.atoms.keys()) and 'C' in list(x.atoms.keys()):
                un_sats.append(x)
    sats = []
    for x in mol_list:
        if x.du is not None and x.du == 0.0:
            if 'H' in list(x.atoms.keys()) and 'C' in list(x.atoms.keys()):
                sats.append(x)

    with open(filename if filename is not None else "unsatpercent.tex", "w") as output:
        output.write(f"{100*len(un_sats)/(len(sats)+len(un_sats)):.0f}\endinput") 

def make_sat_list(mol_list=None, filename=None):
    """
    Makes a .tex file containing the list of molecules from mol_list that are completely
    saturated hydrocarbons

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)   
    filename : str
        The filename for the output snippet (default is 'satlist.tex')
    """ 

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    sats = []
    for x in mol_list:
        # it needs to have a du and for that du to be 0
        if x.du is not None and x.du == 0.0:
            # we only want hydrocarbons here
            if 'H' in list(x.atoms.keys()) and 'C' in list(x.atoms.keys()):
                sats.append(x)

    my_str = ", ".join([r"\ce{" + x.formula + r"}" for x in sats[:-1]])
    my_str += r", and \ce{" + f"{sats[-1].formula}" + r"}\endinput"

    with open(filename if filename is not None else "satlist.tex", "w") as output:
        output.write(my_str)

def make_sat_count(mol_list=None, filename=None):
    """
    Makes a .tex file containing the number of molecules from mol_list that are completely
    saturated hydrocarbons

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)   
    filename : str
        The filename for the output snippet (default is 'nsats.tex')
    """ 

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    sats = []
    for x in mol_list:
        # it needs to have a du and for that du to be 0
        if x.du is not None and x.du == 0.0:
            # we only want hydrocarbons here
            if 'H' in list(x.atoms.keys()) and 'C' in list(x.atoms.keys()):
                sats.append(x)

    with open(filename if filename is not None else "nsats.tex", "w") as output:
        output.write(f"{len(sats)}\endinput")      

def make_sat_percent(mol_list=None, filename=None):
    """
    Makes a .tex file containing the percentage of molecules from mol_list that are completely
    saturated hydrocarbons

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)   
    filename : str
        The filename for the output snippet (default is 'satpercent.tex')
    """ 

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    un_sats = []
    for x in mol_list:
        if x.du is not None and x.du > 0.0:
            if 'H' in list(x.atoms.keys()) and 'C' in list(x.atoms.keys()):
                un_sats.append(x)
    sats = []
    for x in mol_list:
        if x.du is not None and x.du == 0.0:
            if 'H' in list(x.atoms.keys()) and 'C' in list(x.atoms.keys()):
                sats.append(x)

    with open(filename if filename is not None else "satpercent.tex", "w") as output:
        output.write(f"{100*len(sats)/(len(sats)+len(un_sats)):.0f}\endinput")  

def make_sfr_rad_percent(mol_list=None, filename=None):
    """
    Makes a .tex file containing the percentage of molecules from mol_list that are radicals
    and detected in star-forming regions

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)   
    filename : str
        The filename for the output snippet (default is 'sfr_rad_percent.tex')
    """ 

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    mols = []
    for x in mol_list:
        if x.radical:
            if "SFR" in [y.type for y in x.sources]:
                mols.append(x)

    rad_per = 100 * len(mols) / len([x for x in mol_list if "SFR" in [y.type for y in x.sources]])

    with open(filename if filename is not None else "sfr_rad_percent.tex", "w") as output:
        output.write(f"{rad_per:.0f}\endinput")

def make_dark_rad_percent(mol_list=None, filename=None):
    """
    Makes a .tex file containing the percentage of molecules from mol_list that are radicals
    and detected in dark clouds

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)   
    filename : str
        The filename for the output snippet (default is 'dark_rad_percent.tex')
    """ 

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    mols = []
    for x in mol_list:
        if x.radical:
            if "Dark Cloud" in [y.type for y in x.sources]:
                mols.append(x)

    rad_per = 100 * len(mols) / len([x for x in mol_list if "Dark Cloud" in [y.type for y in x.sources]])

    with open(filename if filename is not None else "dark_rad_percent.tex", "w") as output:
        output.write(f"{rad_per:.0f}\endinput")

#############################################################
# 				       PowerPoint Slides 					#
#############################################################

def make_mols_slide(mol_list=None, filename=None):
    """
    Generates a PowerPoint formatted slide containing all detected molecules.

    Parameters
    ----------
    mol_list : list
        A list of molecule objects to use (default is all_molecules)   
    filename : str
        The filename for the output images (default is 'astro_molecules.pptx')
    """

    # If a list wasn't specified, default to all molecules
    if mol_list is None:
        mol_list = all_molecules

    # start working on a presentation
    prs = Presentation()

    # set width and height of slides to the standard 16x9 widescreen PowerPoint format.
    prs.slide_width = Pt(1920)
    prs.slide_height = Pt(1080)

    # make a slide set up for a blank slide; we'll do all the formatting ourselves
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # add the title
    title = slide.shapes.add_textbox(
        left = Inches(0.35),
        top = Inches(0.22),
        width = Inches(18.33),
        height = Inches(1.18),
        )
    
    title.text = f"Known Interstellar Molecules"
    title.text_frame.paragraphs[0].font.size = Pt(64)
    title.text_frame.paragraphs[0].font.name = "Arial"

    # to the top right, add the credit
    credit = slide.shapes.add_textbox(
        left = Inches(22.15),
        top = Inches(0.22),
        width = Pt(300),
        height = Pt(90),
    )

    text_frame = credit.text_frame
    text_frame.clear()
    
    p1 = text_frame.paragraphs[0]
    p1.font.name = 'Arial'
    p1.alignment = PP_ALIGN.RIGHT
    run1 = p1.add_run()
    run1.text = 'Created with '
    run1.font.size = Pt(18)
    run2 = p1.add_run()
    run2.text = 'ASTROMOL '
    run2.font.size = Pt(16)
    run2.font.bold = True
    run3 = p1.add_run()
    run3.text = f'v{version()}'
    run3.font.size = Pt(18)

    p2 = text_frame.add_paragraph()
    p2.font.name = 'Arial'
    p2.alignment = PP_ALIGN.RIGHT
    run4 = p2.add_run()
    run4.text = "bmcguir2.github.io/astromol"
    run4.font.size = Pt(18)

    p3 = text_frame.add_paragraph()
    p3.font.name = 'Arial'
    p3.alignment = PP_ALIGN.RIGHT
    run5 = p3.add_run()
    run5.text = "McGuire 2022 " 
    run5.font.size = Pt(18)
    run6 = p3.add_run()
    run6.text = "ApJS " 
    run6.font.size = Pt(18)
    run6.font.italic = True
    run7 = p3.add_run()
    run7.text = "259, 30" 
    run7.font.size = Pt(18)

    # Make a minifunction to turn a formula into a list of parts
    def _split_formula(formula):
        # split the formula on letters, numbers, and +/- signs
        regex = re.compile(r'(\d+|\s+|\-|\+)')
        return [x for x in regex.split(formula) if x != '']

    # Make a minifunction to turn lists of molecules into text boxes
    def _add_paragraphs(text_box,molecules):
        for mol in molecules:
            p = text_box.text_frame.add_paragraph() # make a new paragraph to hold the molecule
            # run the formula through the ringer, using the table formula if provided
            if mol.table_formula is not None:
                runs = _split_formula(mol.table_formula) 
            else:
                runs = _split_formula(mol.formula)
            p.font.size = Pt(28)
            p.font.name = 'Arial'
            if len(runs) == 1:
                p.text = runs[0]
            else:          
                for i in range(len(runs)):
                    x = runs[i]
                    run = p.add_run()
                    run.text = x
                    # if the first letter is lower case, it's probably 'l' or 'c' and should be italicized
                    if i == 0:
                        if x.islower():
                            run.font.italic = True
                    # we need to hanlde '-' and '+'.
                    if x == '-' or x == '+':
                        # if these are the last character, they're a charge, and should be superscript
                        # if they aren't the last character, it's just a dash and shouldn't be touched
                        if i == len(runs)-1:
                            run.font._element.set('baseline', '30000')
                    # now we deal with numbers
                    elif x.isalpha() is False:
                        # some molecules have isomers designated by starting numbers, these shouldn't be touched
                        if i == 0:
                            pass
                        # for the rest, we need to make sure to subscript them
                        else:
                            run.font._element.set('baseline', '-25000')
                    else:
                        pass 
        _delete_paragraph(text_box.text_frame.paragraphs[0])   

    # to get rid of the empty first paragraph
    def _delete_paragraph(paragraph):
        p = paragraph._p
        parent_element = p.getparent()
        parent_element.remove(p)                

    # a mini class to make it easier to hold meta data for tweaking

    class Group(object):
        def __init__(
            self,
            label=None, # The label that goes at the top of the list
            ncols=None, # The number of columns for an entry
            natoms=None, # The number of atoms to include
            natoms_greater=False, # Set to true to include all molecules with natoms _or greater_.
            label_coords=None, # [left, top, width, height] in inches
            col_coords=None, # [[left, top, width, height],[left, top, width, height]] in inches; only need one entry (but should be double list) for single column
        ):

            self.label = label
            self.ncols = ncols
            self.natoms = natoms
            self.natoms_greater = natoms_greater
            self.label_coords = label_coords
            self.col_coords = col_coords

    # Fill in the info for the arrangement.  Would be nice if it were more automated ...
    
    two_atoms = Group(
        label = "2 Atoms",
        ncols = 2,
        natoms = 2,
        label_coords = [0.35, 1.5, 2.5, 0.8],
        col_coords = [
                        [0.35, 2.1, 1.1, 11.5],
                        [1.55, 2.1, 1.1, 11.5],
                    ],
    )

    three_atoms = Group(
        label = "3 Atoms",
        ncols = 2,
        natoms = 3,
        label_coords = [two_atoms.label_coords[0] + 2.5, 
                        two_atoms.label_coords[1], 
                        two_atoms.label_coords[2], 
                        two_atoms.label_coords[3]],
        col_coords = [
                        [two_atoms.col_coords[0][0] + 2.5, two_atoms.col_coords[0][1], 1.4, 11.],
                        [two_atoms.col_coords[1][0] + 2.7, two_atoms.col_coords[0][1], 1.4, 11.],
                    ],
    )

    four_atoms = Group(
        label = "4 Atoms",
        ncols = 2,
        natoms = 4,
        label_coords = [three_atoms.label_coords[0] + 3.0, 
                        two_atoms.label_coords[1], 
                        two_atoms.label_coords[2], 
                        two_atoms.label_coords[3]],
        col_coords = [
                        [three_atoms.col_coords[0][0] + 3.0, two_atoms.col_coords[0][1], 1.6, 8.],
                        [three_atoms.col_coords[1][0] + 3.2, two_atoms.col_coords[0][1], 1.6, 8.],
                    ],
    )    

    five_atoms = Group(
        label = "5 Atoms",
        ncols = 2,
        natoms = 5,
        label_coords = [four_atoms.label_coords[0] + 3.4, 
                        two_atoms.label_coords[1], 
                        two_atoms.label_coords[2], 
                        two_atoms.label_coords[3]],
        col_coords = [
                        [four_atoms.col_coords[0][0] + 3.4, two_atoms.col_coords[0][1], 1.9, 8.],
                        [four_atoms.col_coords[1][0] + 3.7, two_atoms.col_coords[0][1], 1.9, 8.],
                    ],
    )   

    six_atoms = Group(
        label = "6 Atoms",
        ncols = 1,
        natoms = 6,
        label_coords = [five_atoms.label_coords[0] + 4.0, 
                        two_atoms.label_coords[1], 
                        two_atoms.label_coords[2], 
                        two_atoms.label_coords[3]],
        col_coords = [
                        [five_atoms.col_coords[0][0] + 4.0, two_atoms.col_coords[0][1], 2.5, 10.5],
                    ],
    )   

    seven_atoms = Group(
        label = "7 Atoms",
        ncols = 1,
        natoms = 7,
        label_coords = [six_atoms.label_coords[0] + 2.3, 
                        two_atoms.label_coords[1], 
                        two_atoms.label_coords[2], 
                        two_atoms.label_coords[3]],
        col_coords = [
                        [six_atoms.col_coords[0][0] + 2.3, two_atoms.col_coords[0][1], 2.5, 10.5],
                    ],
    ) 

    eight_atoms = Group(
        label = "8 Atoms",
        ncols = 1,
        natoms = 8,
        label_coords = [seven_atoms.label_coords[0] + 2.5, 
                        two_atoms.label_coords[1], 
                        two_atoms.label_coords[2], 
                        two_atoms.label_coords[3]],
        col_coords = [
                        [seven_atoms.col_coords[0][0] + 2.5, two_atoms.col_coords[0][1], 2.5, 10.5],
                    ],
    )    

    nine_atoms = Group(
        label = "9 Atoms",
        ncols = 2,
        natoms = 9,
        label_coords = [eight_atoms.label_coords[0] + 2.7, 
                        two_atoms.label_coords[1], 
                        two_atoms.label_coords[2], 
                        two_atoms.label_coords[3]],
        col_coords = [
                        [eight_atoms.col_coords[0][0] + 2.7, two_atoms.col_coords[0][1], 3., 4.],
                        [eight_atoms.col_coords[0][0] + 5.7, two_atoms.col_coords[0][1], 3., 4.],
                    ],
    )     

    ten_atoms = Group(
        label = "10 Atoms",
        ncols = 1,
        natoms = 10,
        label_coords = [nine_atoms.label_coords[0], 
                        6.0, 
                        two_atoms.label_coords[2], 
                        two_atoms.label_coords[3]],
        col_coords = [
                        [nine_atoms.col_coords[0][0], 6.45, 3., 3.],
                    ],
    ) 

    eleven_atoms = Group(
        label = "11 Atoms",
        ncols = 1,
        natoms = 11,
        label_coords = [nine_atoms.col_coords[1][0], 
                        6.0, 
                        two_atoms.label_coords[2], 
                        two_atoms.label_coords[3]],
        col_coords = [
                        [nine_atoms.col_coords[1][0], 6.45, 3., 3.],
                    ],
    )  

    twelve_atoms = Group(
        label = "12 Atoms",
        ncols = 1,
        natoms = 12,
        label_coords = [seven_atoms.label_coords[0], 
                        10.6, 
                        two_atoms.label_coords[2], 
                        two_atoms.label_coords[3]],
        col_coords = [
                        [seven_atoms.col_coords[0][0], 11.1, 3., 2.],
                        [eight_atoms.col_coords[0][0], 11.1, 3., 2.],
                    ],
    )   

    more_atoms = Group(
        label = "13+ Atoms",
        ncols = 2,
        natoms = 13,
        natoms_greater = True,
        label_coords = [nine_atoms.label_coords[0], 
                        12.15, 
                        two_atoms.label_coords[2], 
                        two_atoms.label_coords[3]],
        col_coords = [
                        [nine_atoms.col_coords[0][0], 12.6, 3., 2.],
                        [nine_atoms.col_coords[1][0], 12.6, 3., 2.],
                    ],
    )        

    # Loop through and make the items

    groups = [
                two_atoms,
                three_atoms,
                four_atoms,
                five_atoms,
                six_atoms,
                seven_atoms,
                eight_atoms,
                nine_atoms,
                ten_atoms,
                eleven_atoms,
                twelve_atoms,
                more_atoms,
                ]
    my_shapes = {}

    i = 0
    for group in groups:
        my_shapes[i] = slide.shapes.add_textbox(
                            left = Inches(group.label_coords[0]),
                            top = Inches(group.label_coords[1]),
                            width = Inches(group.label_coords[2]),
                            height = Inches(group.label_coords[3]),
                        )

        my_shapes[i].text = group.label
        my_shapes[i].text_frame.paragraphs[0].font.size = Pt(30)
        my_shapes[i].text_frame.paragraphs[0].font.name = 'Arial'
        my_shapes[i].text_frame.paragraphs[0].font.bold = True
        i += 1

        if group.natoms_greater is False:
            molecules = [x for x in mol_list if x.natoms == group.natoms]
        elif group.natoms_greater is True:
            molecules = [x for x in mol_list if x.natoms >= group.natoms]
        for col in range(group.ncols):
            my_shapes[i] = slide.shapes.add_textbox(
                                left = Inches(group.col_coords[col][0]),
                                top = Inches(group.col_coords[col][1]),
                                width = Inches(group.col_coords[col][2]),
                                height = Inches(group.col_coords[col][3]),
                            )

            if col == 0:
                _add_paragraphs(my_shapes[i],molecules[:(math.ceil(len(molecules)/group.ncols))])
            elif col == 1:
                _add_paragraphs(my_shapes[i],molecules[(math.ceil(len(molecules)/group.ncols)):])
            
            i += 1

    # make the number of molecules box

    nbox = slide.shapes.add_shape(
                                    MSO_SHAPE.ROUNDED_RECTANGLE,
                                    left = Inches(6.35),
                                    top = Inches(10.6),
                                    width = Inches(6.0),
                                    height = Inches(1.4),
                                )      

    nbox.fill.solid()
    nbox.fill.fore_color.rgb = RGBColor(194,192,191)
    nbox.line.color.rgb = RGBColor(0,0,0)
    nbox.line.width = Pt(6.)

    nbox.text_frame.paragraphs[0].text = f'{len(mol_list)} Molecules'
    nbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    nbox.text_frame.paragraphs[0].font.size = Pt(45)
    nbox.text_frame.paragraphs[0].font.name = "Arial"
    nbox.text_frame.paragraphs[0].font.bold = True
    nbox.text_frame.paragraphs[0].font.color.rgb = RGBColor(163,31,52)

    as_of = slide.shapes.add_textbox(
        left = Inches(6.35),
        top = Inches(12.2),
        width = Inches(6.),
        height = Inches(1.),
    )    

    as_of.text_frame.paragraphs[0].text = f"Last Updated: {astromol.__updated__.day} {astromol.__updated__.strftime('%b')} {astromol.__updated__.year}"
    as_of.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    as_of.text_frame.paragraphs[0].font.size = Pt(32)
    as_of.text_frame.paragraphs[0].font.name = "Arial"
    as_of.text_frame.paragraphs[0].font.bold = True
            
    # write it out
    if filename is not None:
        prs.save(filename)
    else:
        prs.save('astro_molecules.pptx')


    return